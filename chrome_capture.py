#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Chrome瀏覽器捕獲工具：專門用於通過Chrome瀏覽器擷取網頁中的視頻投影片
"""

import os
import time
import json
import argparse
import tkinter as tk
from tkinter import messagebox, Scale, filedialog, ttk
import numpy as np
import cv2
import pyautogui
from datetime import datetime
from PIL import Image, ImageTk
from skimage.metrics import structural_similarity as ssim
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
import subprocess
import sys
import threading
import traceback

# 導入 PPT 生成函數
try:
    from video_audio_processor import generate_ppt_from_images
except ImportError:
    # 如果找不到該模塊，這裡添加 generate_ppt_from_images 函數的實現
    def generate_ppt_from_images(image_folder, output_file=None, title="視頻捕獲的幻燈片", 
                              image_files=None, slide_ratio="16:9"):
        """
        將圖片文件夾轉換為 PowerPoint 文件
        
        參數:
            image_folder: 包含幻燈片圖片的文件夾
            output_file: 輸出的 PowerPoint 文件路徑，默認為文件夾名+.pptx
            title: 演示文稿標題
            image_files: 已排序的圖片文件列表，若為 None 則自動從文件夾獲取
            slide_ratio: 幻燈片比例，支持 "16:9" 或 "4:3"
            
        返回:
            success: 是否成功
            output_file: 輸出的文件路徑，或者錯誤訊息
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            # 處理默認輸出文件
            if not output_file:
                folder_name = os.path.basename(image_folder)
                output_file = os.path.join(
                    os.path.dirname(image_folder),
                    f"{folder_name}.pptx"
                )
            
            # 確保輸出路徑有效
            if not output_file or output_file.strip() == "" or output_file.strip() == ".pptx":
                # 使用文件夾名稱作為檔案名
                folder_name = os.path.basename(image_folder)
                if not folder_name or folder_name.strip() == "":
                    folder_name = "投影片"  # 默認檔案名
                    
                output_file = os.path.join(
                    os.path.dirname(image_folder),
                    f"{folder_name}.pptx"
                )
            
            # 添加副檔名
            if not output_file.lower().endswith(".pptx"):
                output_file += ".pptx"
            
            # 如果沒有傳入圖片列表，從文件夾獲取
            if not image_files:
                image_files = []
                for filename in os.listdir(image_folder):
                    if filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                        image_files.append(os.path.join(image_folder, filename))
                image_files.sort()  # 按文件名排序
                
            if not image_files:
                return False, "未找到圖片文件"
            
            # 創建演示文稿對象
            prs = Presentation()
            
            # 設置幻燈片尺寸
            if slide_ratio == "16:9":
                # 16:9比例的尺寸 (寬度10英寸, 高度5.625英寸)
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(5.625)
            else:  # 4:3
                # 4:3比例的尺寸 (寬度10英寸, 高度7.5英寸)
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
                
            # 添加標題幻燈片
            title_slide_layout = prs.slide_layouts[0]  # 第一個是標題幻燈片
            slide = prs.slides.add_slide(title_slide_layout)
            
            # 設置標題
            title_shape = slide.shapes.title
            title_shape.text = title
            
            # 設置副標題
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = f"自動生成於 {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            
            # 添加投影片
            for i, img_path in enumerate(image_files):
                # 添加空白投影片
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白佈局
                
                # 獲取圖像尺寸
                img = Image.open(img_path)
                width, height = img.size
                
                # 計算投影片中的位置和大小
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # 保持寬高比例
                ratio = min(slide_width / width, slide_height / height)
                
                # 添加圖片
                left = (slide_width - width * ratio) / 2
                top = (slide_height - height * ratio) / 2
                slide.shapes.add_picture(
                    img_path, left, top, 
                    width=width * ratio, height=height * ratio
                )
            
            # 保存為PowerPoint
            prs.save(output_file)
            return True, output_file
            
        except Exception as e:
            error_msg = str(e)
            return False, f"生成PowerPoint時出錯: {error_msg}"

class ChromeCapture:
    def __init__(self, root=None):
        # 創建主窗口
        if not root:
            self.root = tk.Tk()
            self.root.title("Chrome瀏覽器投影片捕獲工具")
            self.root.geometry("1000x700")
        else:
            self.root = root
        
        # 加載配置
        self.load_config()
        
        # 初始化图片缓存
        self.photo_cache = []
        self.current_photo = None
        
        # 設置自動停止相關屬性
        self.inactivity_timeout = 2 * 60  # 預設2分鐘無變化自動停止（秒）
        self.auto_stop_enabled = True  # 是否啟用自動停止功能
        self.last_change_time = time.time()  # 記錄最後一次檢測到變化的時間
        
        # Slide 製作相關變數
        self.slide_make_status_var = tk.StringVar()
        self.slide_make_status_var.set("就緒")
        self.preview_images = []  # 保存預覽圖像引用
        self.sort_option_var = tk.StringVar(value="name")  # 預設按文件名排序
        self.slide_ratio_var = tk.StringVar(value="16:9")  # 預設寬螢幕比例
        
        # 設置 UI
        self.setup_ui()
        
        # 設置初始狀態
        self.browser = None
        self.frame_size = None
        self.display_size = None
        self.roi = None
        self.roi_start = None
        self.roi_end = None
        self.last_mouse_pos = (0, 0)  # 記錄滑鼠最後位置
        self.roi_selecting = False    # 是否正在選擇ROI
        self.capture_running = False
        self.status_var = tk.StringVar()
        self.status_var.set("就緒")
        self.canvas_offset = (0, 0)
        
        # 確保輸出目錄存在
        if not os.path.exists(self.output_folder):
            os.makedirs(self.output_folder)
        
        # 事件綁定
        self.root.protocol("WM_DELETE_WINDOW", self.on_closing)
    
    def load_config(self):
        """從.cursor.json加載配置"""
        try:
            if os.path.exists('.cursor.json'):
                with open('.cursor.json', 'r', encoding='utf-8') as f:
                    config = json.load(f)
                    
                # 從配置中獲取Chrome捕獲設置
                if ('customSettings' in config and 
                        'chromeCapture' in config['customSettings']):
                    chrome_settings = config['customSettings']['chromeCapture']
                    self.threshold = chrome_settings.get(
                        'defaultThreshold', 0.95)
                    self.interval = chrome_settings.get(
                        'defaultInterval', 0.5)
                    self.output_format = chrome_settings.get(
                        'defaultOutputFormat', 'pptx')
                    self.auto_save = chrome_settings.get(
                        'autoSaveSlides', True)
                    self.chrome_path = chrome_settings.get(
                        'browserExecutablePath', '')
                    self.user_data_dir = chrome_settings.get(
                        'userDataDir', '')
                else:
                    self.threshold = 0.95
                    self.interval = 0.5
                    self.output_format = 'pptx'
                    self.auto_save = True
                    self.chrome_path = ''
                    self.user_data_dir = ''
                
                # 從配置中獲取捕獲設置
                if ('customSettings' in config and 
                        'captureSettings' in config['customSettings']):
                    capture_settings = (
                        config['customSettings']['captureSettings']
                    )
                    self.screenshot_method = capture_settings.get(
                        'preferredScreenShotMethod', 'pyautogui')
                    self.fallback_method = capture_settings.get(
                        'fallbackMethod', 'opencv')
                    self.similarity_algorithm = capture_settings.get(
                        'similarityAlgorithm', 'ssim')
                else:
                    self.screenshot_method = 'pyautogui'
                    self.fallback_method = 'opencv'
                    self.similarity_algorithm = 'ssim'
            else:
                # 默認值
                self.threshold = 0.95
                self.interval = 0.5
                self.output_format = 'pptx'
                self.auto_save = True
                self.chrome_path = ''
                self.user_data_dir = ''
                self.screenshot_method = 'pyautogui'
                self.fallback_method = 'opencv'
                self.similarity_algorithm = 'ssim'
        except Exception as e:
            print(f"加載配置時出錯: {str(e)}")
            # 默認值
            self.threshold = 0.95
            self.interval = 0.5
            self.output_format = 'pptx'
            self.auto_save = True
            self.chrome_path = ''
            self.user_data_dir = ''
            self.screenshot_method = 'pyautogui'
            self.fallback_method = 'opencv'
            self.similarity_algorithm = 'ssim'
            
        # 其他設置
        self.output_folder = "slides"
        self.output_file = f"slides.{self.output_format}"
    
    def setup_ui(self):
        # 創建標籤頁系統
        self.tab_control = ttk.Notebook(self.root)
        
        # 創建主捕獲標籤頁
        self.main_tab = ttk.Frame(self.tab_control)
        self.tab_control.add(self.main_tab, text="投影片捕獲")
        
        # 創建 Slide 製作標籤頁
        self.slide_make_tab = ttk.Frame(self.tab_control)
        self.tab_control.add(self.slide_make_tab, text="Slide 製作")
        
        # 顯示標籤頁
        self.tab_control.pack(expand=1, fill=tk.BOTH, padx=5, pady=5)
        
        # 設置主捕獲標籤頁 UI
        self.setup_main_capture_ui()
        
        # 設置 Slide 製作標籤頁 UI
        self.setup_slide_make_ui()
    
    def setup_main_capture_ui(self):
        """設置主捕獲標籤頁 UI"""
        # 上方控制區域
        control_frame = tk.Frame(self.main_tab)
        control_frame.pack(fill=tk.X, padx=10, pady=10)
        
        # 第一行 - 網址輸入
        row0 = tk.Frame(control_frame)
        row0.pack(fill=tk.X, pady=5)
        
        tk.Label(row0, text="網址:").pack(side=tk.LEFT, padx=5)
        self.url_entry = tk.Entry(row0, width=50)
        self.url_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        # 為網址輸入框創建右鍵菜單
        self.create_entry_context_menu(self.url_entry)
        
        self.go_btn = tk.Button(
            row0, text="打開瀏覽器", command=self.launch_browser, 
            bg="#2196F3", fg="white", width=12
        )
        self.go_btn.pack(side=tk.LEFT, padx=5)
        
        # 瀏覽器選擇區域按鈕 (初始狀態為隱藏)
        self.browser_select_btn = tk.Button(
            row0, text="在瀏覽器中選擇", command=self.select_area_in_browser,
            bg="#FF9800", fg="white", width=12
        )
        # 開始時不顯示，等瀏覽器啟動後再顯示
        
        # 第二行 - 設置
        row1 = tk.Frame(control_frame)
        row1.pack(fill=tk.X, pady=5)
        
        tk.Label(row1, text="相似度閾值:").pack(side=tk.LEFT, padx=5)
        self.threshold_scale = Scale(
            row1, from_=0.5, to=1.0, resolution=0.01, 
            orient=tk.HORIZONTAL, length=200
        )
        self.threshold_scale.set(self.threshold)
        self.threshold_scale.pack(side=tk.LEFT, padx=5)
        
        tk.Label(row1, text="檢測間隔(秒):").pack(side=tk.LEFT, padx=5)
        self.interval_scale = Scale(
            row1, from_=0.1, to=5.0, resolution=0.1, 
            orient=tk.HORIZONTAL, length=200
        )
        self.interval_scale.set(self.interval)
        self.interval_scale.pack(side=tk.LEFT, padx=5)
        
        # 新增設置 - 無變化自動停止
        row1_5 = tk.Frame(control_frame)
        row1_5.pack(fill=tk.X, pady=5)
        
        tk.Label(row1_5, text="無變化自動停止(分鐘):").pack(side=tk.LEFT, padx=5)
        self.timeout_scale = Scale(
            row1_5, from_=0, to=30, resolution=1, 
            orient=tk.HORIZONTAL, length=200
        )
        self.timeout_scale.set(self.inactivity_timeout // 60)  # 轉換秒為分鐘
        self.timeout_scale.pack(side=tk.LEFT, padx=5)
        
        # 添加自動停止功能啟用/禁用選項
        self.auto_stop_var = tk.BooleanVar(value=self.auto_stop_enabled)
        self.auto_stop_checkbox = tk.Checkbutton(
            row1_5, text="啟用自動停止", variable=self.auto_stop_var,
            command=self.toggle_auto_stop
        )
        self.auto_stop_checkbox.pack(side=tk.LEFT, padx=5)
        
        # 第三行 - 按鈕
        row2 = tk.Frame(control_frame)
        row2.pack(fill=tk.X, pady=5)
        
        self.start_btn = tk.Button(
            row2, text="開始捕獲", command=self.start_capture, 
            bg="#4CAF50", fg="white", width=10
        )
        self.start_btn.pack(side=tk.LEFT, padx=5)
        
        self.pause_btn = tk.Button(
            row2, text="暫停", command=self.pause_capture, 
            state=tk.DISABLED, width=10
        )
        self.pause_btn.pack(side=tk.LEFT, padx=5)
        
        self.stop_btn = tk.Button(
            row2, text="停止", command=self.stop_capture, 
            state=tk.DISABLED, width=10
        )
        self.stop_btn.pack(side=tk.LEFT, padx=5)
        
        self.reset_roi_btn = tk.Button(
            row2, text="重設區域", command=self.reset_roi, width=10
        )
        self.reset_roi_btn.pack(side=tk.LEFT, padx=5)
        
        self.generate_btn = tk.Button(
            row2, text="生成PPT", command=self.generate_ppt, width=10
        )
        self.generate_btn.pack(side=tk.LEFT, padx=5)
        
        # 右側狀態訊息
        self.status_var = tk.StringVar(value="就緒")
        status_label = tk.Label(
            row2, textvariable=self.status_var, 
            bd=1, relief=tk.SUNKEN, anchor=tk.W
        )
        status_label.pack(side=tk.RIGHT, padx=5, fill=tk.X, expand=True)
        
        # 預覽區域
        self.canvas_frame = tk.Frame(self.main_tab, bg="black")
        self.canvas_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        self.canvas = tk.Canvas(self.canvas_frame, bg="black")
        self.canvas.pack(fill=tk.BOTH, expand=True)
        
        # 綁定滑鼠事件
        self.canvas.bind("<ButtonPress-1>", self.on_mouse_down)
        self.canvas.bind("<B1-Motion>", self.on_mouse_move)
        self.canvas.bind("<ButtonRelease-1>", self.on_mouse_up)
        
        # 日誌區域
        log_frame = tk.Frame(self.main_tab)
        log_frame.pack(fill=tk.X, padx=10, pady=5)
        
        tk.Label(log_frame, text="執行日誌:").pack(anchor=tk.W)
        self.log_text = tk.Text(log_frame, height=6, width=50)
        self.log_text.pack(fill=tk.X, expand=True)
        
        scrollbar = tk.Scrollbar(self.log_text)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.log_text.config(yscrollcommand=scrollbar.set)
        scrollbar.config(command=self.log_text.yview)
        
        self.log("程序已啟動")
        self.log("請輸入網址並打開瀏覽器，然後框選要監控的投影片區域")
    
    def log(self, message):
        timestamp = datetime.now().strftime("%H:%M:%S")
        self.log_text.insert(tk.END, f"[{timestamp}] {message}\n")
        self.log_text.see(tk.END)
    
    def launch_browser(self):
        url = self.url_entry.get().strip()
        if not url:
            self.log("請輸入有效的網址")
            return
        
        # 檢查 URL 格式
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url
            self.url_entry.delete(0, tk.END)
            self.url_entry.insert(0, url)
        
        try:
            # 啟動 Chrome
            if not self.browser:
                chrome_options = Options()
                chrome_options.add_argument("--start-maximized")
                self.browser = webdriver.Chrome(options=chrome_options)
            
            # 導航到網址
            self.browser.get(url)
            
            # 顯示瀏覽器預覽
            self.show_browser_preview()
            
            # 顯示瀏覽器選擇按鈕
            self.browser_select_btn.pack(side=tk.LEFT, padx=5)
            
            self.log("瀏覽器已啟動，請點擊「在瀏覽器中選擇」進行精確選擇")
        except Exception as e:
            self.log(f"啟動瀏覽器時出錯: {str(e)}")
    
    def select_area_in_browser(self):
        """在瀏覽器中直接選擇擷取區域"""
        try:
            # 這裡暫時導入，是為了在使用時才檢查是否安裝
            import pynput
        except ImportError:
            self.log("未安裝必要的套件: pynput，正在嘗試安裝...")
            try:
                subprocess.check_call(
                    [sys.executable, "-m", "pip", "install", "pynput"]
                )
                self.log("安裝成功，請重新點擊「在瀏覽器中選擇」按鈕")
                return
            except Exception as e:
                self.log(f"安裝失敗: {str(e)}，請手動安裝 pynput: pip install pynput")
                return
        
        if not self.browser:
            self.log("請先啟動瀏覽器")
            return
        
        # 最小化應用程式視窗
        self.log("正在準備瀏覽器區域選擇模式...")
        self.log("請在瀏覽器視窗上按下滑鼠左鍵並拖曳來選擇區域")
        self.log("完成選擇後，應用程式將自動恢復")
        
        # 給使用者時間準備
        self.root.after(1000, self._start_browser_selection)
    
    def _start_browser_selection(self):
        """開始瀏覽器區域選擇的實際過程"""
        # 最小化應用程式視窗
        self.root.iconify()
        
        # 確保瀏覽器視窗在前景
        self.browser.maximize_window()
        
        try:
            from pynput.mouse import Listener, Button
            
            # 用於儲存選擇的座標
            selection_coords = {
                'start_x': 0, 'start_y': 0, 
                'end_x': 0, 'end_y': 0, 
                'selecting': False
            }
            
            def on_click(x, y, button, pressed):
                if button == Button.left:
                    if pressed:
                        # 開始選擇
                        selection_coords['start_x'] = x
                        selection_coords['start_y'] = y
                        selection_coords['selecting'] = True
                    else:
                        # 結束選擇
                        if selection_coords['selecting']:
                            selection_coords['end_x'] = x
                            selection_coords['end_y'] = y
                            selection_coords['selecting'] = False
                            return False  # 停止監聽
                elif button == Button.right:
                    # 右鍵取消選擇
                    return False
                return True  # 繼續監聽
            
            # 啟動監聽器
            with Listener(on_click=on_click) as listener:
                listener.join()
            
            # 處理選擇結果
            self.root.after(
                500, 
                lambda: self._process_browser_selection(selection_coords)
            )
        except Exception as e:
            self.log(f"選擇過程中發生錯誤: {str(e)}")
            self.root.deiconify()  # 恢復視窗
    
    def _process_browser_selection(self, coords):
        """處理瀏覽器區域選擇的結果"""
        # 恢復應用程式視窗
        self.root.deiconify()
        
        # 確保選擇有效（開始點與結束點不同）
        if (coords['start_x'] == coords['end_x'] or 
                coords['start_y'] == coords['end_y']):
            self.log("選擇無效，請重新選擇有效的區域")
            return
        
        # 確保座標是遞增的
        x1 = min(coords['start_x'], coords['end_x'])
        y1 = min(coords['start_y'], coords['end_y'])
        x2 = max(coords['start_x'], coords['end_x'])
        y2 = max(coords['start_y'], coords['end_y'])
        
        # 設定 ROI
        self.roi = (x1, y1, x2, y2)
        
        # 更新顯示
        self.log(f"已設定選擇區域: ({x1}, {y1}) - ({x2}, {y2})")
        
        # 重新顯示預覽
        self.show_browser_preview()
    
    def show_browser_preview(self):
        """顯示瀏覽器預覽圖"""
        if self.browser is None:
            self.status_var.set("瀏覽器未啟動")
            return
        
        try:
            # 確保應用程序窗口在前景
            self.root.deiconify()
            self.root.lift()
            self.root.focus_force()
            
            # 等待 UI 更新
            self.root.update()
            self.root.after(200, self._take_browser_screenshot)
            
        except Exception as e:
            error_msg = f"無法獲取瀏覽器預覽: {str(e)}"
            print(error_msg)
            self.status_var.set(error_msg)
            messagebox.showerror("錯誤", error_msg)
            import traceback
            traceback.print_exc()
    
    def _take_browser_screenshot(self):
        """截取瀏覽器螢幕畫面"""
        try:
            # 截取全螢幕
            self.log("正在擷取瀏覽器畫面...")
            screen = pyautogui.screenshot()
            screen_np = np.array(screen)
            screen_bgr = cv2.cvtColor(screen_np, cv2.COLOR_RGB2BGR)
            
            # 設置實際幀大小
            self.frame_size = (screen_bgr.shape[1], screen_bgr.shape[0])
            
            # 顯示到UI
            self.display_frame(screen_bgr)
            self.status_var.set("請在瀏覽器中框選投影片區域")
            
            # 將browser_select_btn按鈕顯示出來
            if not self.browser_select_btn.winfo_ismapped():
                self.browser_select_btn.pack(side=tk.LEFT, padx=5)
            
            # 記錄日誌
            self.log(f"已擷取瀏覽器畫面 ({self.frame_size[0]}x{self.frame_size[1]})")
            
        except Exception as e:
            error_msg = f"擷取瀏覽器畫面時出錯: {str(e)}"
            print(error_msg)
            self.status_var.set(error_msg)
            self.log(error_msg)
            import traceback
            traceback.print_exc()
    
    def on_mouse_down(self, event):
        # 檢查是否處於可選擇狀態，不管是否正在運行都可以重設選擇
        self.roi_selecting = True
        self.roi_start = (event.x, event.y)
        self.canvas.delete("roi")
    
    def on_mouse_move(self, event):
        if self.roi_selecting:
            self.roi_end = (event.x, event.y)
            self.canvas.delete("roi_temp")
            self.canvas.create_rectangle(
                self.roi_start[0], self.roi_start[1],
                self.roi_end[0], self.roi_end[1],
                outline="red", width=2, tags="roi_temp"
            )
    
    def on_mouse_up(self, event):
        if self.roi_selecting:
            self.roi_selecting = False
            self.roi_end = (event.x, event.y)
            
            # 計算 ROI
            x1 = min(self.roi_start[0], self.roi_end[0])
            y1 = min(self.roi_start[1], self.roi_end[1])
            x2 = max(self.roi_start[0], self.roi_end[0])
            y2 = max(self.roi_start[1], self.roi_end[1])
            
            # 確保區域有足夠大小
            if x2 - x1 < 5 or y2 - y1 < 5:
                messagebox.showerror("錯誤", "選擇的區域太小，請重新選擇")
                self.canvas.delete("roi_temp")
                return
            
            # 將畫布上的座標轉換為螢幕實際座標
            # 首先獲取圖像的實際位置和大小
            display_width, display_height = self.display_size
            offset_x, offset_y = self.canvas_offset
            
            # 計算比例
            screen_width, screen_height = self.frame_size
            scale_x = screen_width / display_width
            scale_y = screen_height / display_height
            
            # 調整座標以考慮偏移和縮放
            screen_x1 = int((x1 - offset_x) * scale_x)
            screen_y1 = int((y1 - offset_y) * scale_y)
            screen_x2 = int((x2 - offset_x) * scale_x)
            screen_y2 = int((y2 - offset_y) * scale_y)
            
            # 確保座標在範圍內
            screen_x1 = max(0, min(screen_x1, screen_width))
            screen_y1 = max(0, min(screen_y1, screen_height))
            screen_x2 = max(0, min(screen_x2, screen_width))
            screen_y2 = max(0, min(screen_y2, screen_height))
            
            # 保存實際像素座標為ROI
            self.roi = (screen_x1, screen_y1, screen_x2, screen_y2)
            
            # 在畫布上保存選擇區域
            self.canvas.delete("roi_temp")
            self.canvas.create_rectangle(
                x1, y1, x2, y2,
                outline="green", width=2, tags="roi_canvas"
            )
            
            # 顯示實際像素座標
            rx1, ry1, rx2, ry2 = self.roi
            self.log(f"已選擇區域: ({rx1},{ry1}) - ({rx2},{ry2}) [實際像素]")
            
            # 繪製實際ROI參考框（按照實際比例）
            self.draw_actual_roi()
    
    def draw_actual_roi(self):
        """按照實際比例繪製ROI框"""
        if not self.roi or not hasattr(self, 'display_size') or not hasattr(self, 'frame_size'):
            return
        
        # 清除舊的ROI參考
        self.canvas.delete("roi")
        
        # 獲取ROI座標
        x1, y1, x2, y2 = self.roi
        
        # 計算畫布中的位置
        display_width, display_height = self.display_size
        screen_width, screen_height = self.frame_size
        offset_x, offset_y = self.canvas_offset
        
        # 計算比例
        scale_x = display_width / screen_width
        scale_y = display_height / screen_height
        
        # 計算畫布上的位置
        canvas_x1 = offset_x + int(x1 * scale_x)
        canvas_y1 = offset_y + int(y1 * scale_y)
        canvas_x2 = offset_x + int(x2 * scale_x)
        canvas_y2 = offset_y + int(y2 * scale_y)
        
        # 繪製ROI框
        self.canvas.create_rectangle(
            canvas_x1, canvas_y1, canvas_x2, canvas_y2,
            outline="red", width=2, tags="roi"
        )
        
        self.canvas.update()
    
    def reset_roi(self):
        """重設選擇的區域"""
        self.roi = None
        self.roi_start = None
        self.roi_end = None
        self.canvas.delete("roi")
        self.canvas.delete("roi_canvas")
        self.log("已重設選擇區域")
        
        try:
            # 重新顯示當前螢幕以便選擇新區域
            if self.browser:
                # 優先使用瀏覽器預覽
                self.root.after(100, self.show_browser_preview)
            else:
                # 如果沒有瀏覽器，直接截取全螢幕
                self.log("正在擷取螢幕畫面...")
                
                # 確保UI更新並可見
                self.root.update()
                self.root.lift()
                self.root.focus_force()
                
                # 等待UI完全顯示
                self.root.after(300, self._take_screenshot_and_display)
                
        except Exception as e:
            error_msg = f"重設區域時出錯: {str(e)}"
            print(error_msg)
            self.log(error_msg)
            import traceback
            traceback.print_exc()
    
    def _take_screenshot_and_display(self):
        """截取螢幕畫面並顯示"""
        try:
            # 截取全螢幕
            screenshot = pyautogui.screenshot()
            frame = np.array(screenshot)
            frame = cv2.cvtColor(frame, cv2.COLOR_RGB2BGR)
            
            # 設置frame_size
            self.frame_size = (frame.shape[1], frame.shape[0])
            
            # 直接呼叫display_frame顯示畫面
            self.display_frame(frame)
            
            # 輸出日誌
            self.log(f"已截取螢幕畫面 ({frame.shape[1]}x{frame.shape[0]})")
            self.status_var.set("請選擇要監控的區域")
        except Exception as e:
            error_msg = f"截取螢幕畫面時出錯: {str(e)}"
            print(error_msg)
            self.log(error_msg)
            import traceback
            traceback.print_exc()
    
    def display_frame(self, frame):
        """顯示影像幀"""
        try:
            # 確保canvas已準備好
            self.root.update_idletasks()
            
            # 轉換BGR為RGB
            rgb_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            
            # 調整大小以適應畫布
            canvas_width = self.canvas.winfo_width()
            canvas_height = self.canvas.winfo_height()
            
            # 如果畫布尺寸不合理（可能是尚未更新），使用默認尺寸
            if canvas_width < 10 or canvas_height < 10:
                canvas_width = 800
                canvas_height = 500
            
            # 保持縱橫比例
            img_height, img_width = rgb_frame.shape[:2]
            aspect_ratio = img_width / img_height
            
            if canvas_width / canvas_height > aspect_ratio:
                new_height = canvas_height
                new_width = int(new_height * aspect_ratio)
            else:
                new_width = canvas_width
                new_height = int(new_width / aspect_ratio)
            
            # 確保尺寸至少為1像素
            new_width = max(1, new_width)
            new_height = max(1, new_height)
            
            # 設置frame_size，用於ROI選擇
            self.frame_size = (img_width, img_height)
                
            # 儲存顯示尺寸，用於ROI選擇
            self.display_size = (new_width, new_height)
            
            # 調整大小
            resized_frame = cv2.resize(
                rgb_frame, (new_width, new_height), 
                interpolation=cv2.INTER_AREA
            )
            
            # 轉換為PIL圖像
            img = Image.fromarray(resized_frame)
            
            # 清除畫布
            self.canvas.delete("all")
            
            # 計算居中偏移量
            offset_x = (canvas_width - new_width) // 2
            offset_y = (canvas_height - new_height) // 2
            
            # 保存當前偏移量以供座標轉換使用
            self.canvas_offset = (offset_x, offset_y)
            
            # 創建新的Photo對象
            photo = ImageTk.PhotoImage(image=img)
            
            # 保存到緩存中避免被垃圾回收
            self.photo_cache.append(photo)
            
            # 限制緩存大小
            if len(self.photo_cache) > 10:
                self.photo_cache = self.photo_cache[-5:]
            
            # 更新當前顯示的圖片
            self.current_photo = photo
            
            # 顯示圖像
            self.canvas.create_image(
                offset_x, offset_y,
                image=self.current_photo, anchor=tk.NW, tags="frame"
            )
            
            # 強制更新畫布
            self.canvas.update()
            
            # 如果有ROI，重新繪製
            if hasattr(self, 'roi') and self.roi:
                self.draw_actual_roi()
                
        except Exception as e:
            error_msg = f"顯示幀時出錯: {str(e)}"
            print(error_msg)
            self.status_var.set(error_msg)
            traceback.print_exc()
    
    def display_current_roi(self, frame):
        """顯示當前ROI參考框"""
        if self.roi is None or self.display_size is None:
            return
            
        # 獲取原始框架尺寸
        img_height, img_width = frame.shape[:2]
        
        # 顯示的尺寸
        display_width, display_height = self.display_size
        
        # 原始ROI座標
        x, y, w, h = self.roi
        
        # 計算顯示畫面上的ROI座標
        display_x = int(x * display_width / img_width)
        display_y = int(y * display_height / img_height)
        display_w = int(w * display_width / img_width)
        display_h = int(h * display_height / img_height)
        
        # 獲取畫布中心坐標以正確放置 ROI 框
        canvas_width = self.canvas.winfo_width()
        canvas_height = self.canvas.winfo_height()
        
        # 計算圖像開始的位置
        img_x = canvas_width // 2 - display_width // 2
        img_y = canvas_height // 2 - display_height // 2
        
        # 計算調整後的 ROI 座標
        adjusted_x = img_x + display_x
        adjusted_y = img_y + display_y
        
        # 在畫布上繪製矩形框
        try:
            self.canvas.create_rectangle(
                adjusted_x, adjusted_y, 
                adjusted_x + display_w, adjusted_y + display_h,
                outline="red", width=2, tags="roi"
            )
            self.canvas.update()  # 確保立即顯示
        except Exception as e:
            print(f"顯示ROI框時出錯: {str(e)}")
    
    def display_selection_box(self):
        """顯示當前選擇框"""
        if not self.roi_start or not hasattr(self, 'canvas_offset'):
            return
            
        x1, y1 = self.roi_start
        
        # 獲取當前位置
        if self.roi_end:
            x2, y2 = self.roi_end
        else:
            # 使用滑鼠當前位置
            x2, y2 = self.last_mouse_pos if hasattr(self, 'last_mouse_pos') else (x1, y1)
            
        offset_x, offset_y = self.canvas_offset
        
        # 移除舊的選擇框
        self.canvas.delete("selection")
        
        # 繪製新的選擇框
        self.canvas.create_rectangle(
            x1 + offset_x, y1 + offset_y, 
            x2 + offset_x, y2 + offset_y,
            outline="green", width=2, tags="selection"
        )
    
    def start_capture(self):
        """開始捕獲投影片"""
        if not self.browser:
            messagebox.showerror("錯誤", "請先啟動瀏覽器")
            return
        
        if not self.roi:
            messagebox.showerror("錯誤", "請先選擇要監控的區域")
            return
        
        # 獲取當前設置
        self.threshold = self.threshold_scale.get()
        self.interval = self.interval_scale.get()
        self.auto_stop_enabled = self.auto_stop_var.get()
        # 將分鐘轉換為秒
        self.inactivity_timeout = self.timeout_scale.get() * 60  
        # 重設最後變化時間
        self.last_change_time = time.time()  
        
        # 確認 ROI 座標在當前螢幕截圖上的有效性
        try:
            # 獲取最新的螢幕截圖
            screenshot = pyautogui.screenshot()
            frame = np.array(screenshot)
            frame = cv2.cvtColor(frame, cv2.COLOR_RGB2BGR)
            
            # 更新螢幕尺寸
            screen_w, screen_h = frame.shape[1], frame.shape[0]
            self.frame_size = (screen_w, screen_h)
            
            # 確保 ROI 區域在螢幕範圍內
            x1, y1, x2, y2 = self.roi
            
            if x1 < 0 or y1 < 0 or x2 > screen_w or y2 > screen_h:
                self.log("警告：選擇的區域超出螢幕範圍，已自動調整")
                x1 = max(0, min(x1, screen_w - 1))
                y1 = max(0, min(y1, screen_h - 1))
                x2 = max(0, min(x2, screen_w - 1))
                y2 = max(0, min(y2, screen_h - 1))
                
                # 確保區域有足夠大小
                if x2 - x1 < 5 or y2 - y1 < 5:
                    messagebox.showerror("錯誤", "選擇的區域太小，請重新選擇")
                    return
                
                self.roi = (x1, y1, x2, y2)
                
                # 更新顯示的ROI框
                self.draw_actual_roi()
            
            # 初始化捕獲所需的變數
            self.slides = []
            self.previous_frame = None
            self.slide_count = 0  # 重設投影片計數
            self.is_paused = False
            
            # 顯示最新的一幀帶 ROI 的預覽
            preview_frame = frame.copy()
            
            # 在UI中顯示
            self.display_frame(preview_frame)
            self.draw_actual_roi()  # 再次繪製ROI確保顯示正確
            
            # 記錄實際捕獲區域
            self.log(f"實際捕獲區域: ({x1},{y1}) - ({x2},{y2})")
            
        except Exception as e:
            import traceback
            error_msg = f"獲取螢幕截圖時出錯: {str(e)}"
            self.log(error_msg)
            messagebox.showerror("錯誤", error_msg)
            print(traceback.format_exc())
            return
        
        self.capture_running = True
        
        # 更新按鈕狀態
        self.start_btn.config(state=tk.DISABLED)
        self.pause_btn.config(state=tk.NORMAL)
        self.stop_btn.config(state=tk.NORMAL)
        self.go_btn.config(state=tk.DISABLED)
        
        # 啟動捕獲線程
        self.capture_thread = threading.Thread(target=self.capture_loop)
        self.capture_thread.daemon = True
        self.capture_thread.start()
        
        msg = (f"開始捕獲投影片 (閾值: {self.threshold}, "
               f"間隔: {self.interval}秒")
        if self.auto_stop_enabled and self.inactivity_timeout > 0:
            msg += (f", 無變化自動停止: "
                    f"{self.inactivity_timeout//60}分鐘)")
        else:
            msg += ", 自動停止功能已禁用)"
        self.log(msg)
        self.status_var.set("正在捕獲")
        
    def capture_loop(self):
        """捕獲循環"""
        while self.capture_running:
            if not hasattr(self, 'is_paused') or not self.is_paused:
                try:
                    # 獲取螢幕截圖
                    screenshot = pyautogui.screenshot()
                    
                    # 轉換為 OpenCV 格式
                    frame = np.array(screenshot)
                    frame = cv2.cvtColor(frame, cv2.COLOR_RGB2BGR)
                    
                    # 顯示帶 ROI 的預覽
                    preview_frame = frame.copy()
                    
                    # 在UI中顯示
                    self.root.after(
                        10, lambda f=preview_frame: self.display_frame(f)
                    )
                    # 使用新方法繪製ROI
                    self.root.after(20, self.draw_actual_roi)
                    
                    # 剪切 ROI 區域
                    if self.roi:
                        x1, y1, x2, y2 = self.roi
                        # 確保座標不超出範圍
                        h, w = frame.shape[:2]
                        x1 = max(0, min(x1, w-1))
                        y1 = max(0, min(y1, h-1))
                        x2 = max(0, min(x2, w-1))
                        y2 = max(0, min(y2, h-1))
                        roi_frame = frame[y1:y2, x1:x2]
                        
                        # 檢查是否需要保存
                        if self.previous_frame is None:
                            # 第一幀直接保存
                            self.save_slide(roi_frame)
                            self.previous_frame = roi_frame
                            self.last_change_time = time.time()  # 初始化最後變化時間
                        else:
                            # 計算相似度
                            if self.similarity_algorithm == 'ssim':
                                # 將兩幀調整為相同大小
                                if roi_frame.shape != self.previous_frame.shape:
                                    roi_frame = cv2.resize(
                                        roi_frame, 
                                        (self.previous_frame.shape[1], 
                                         self.previous_frame.shape[0])
                                    )
                                    
                                # 轉換為灰度圖
                                gray_current = cv2.cvtColor(
                                    roi_frame, cv2.COLOR_BGR2GRAY
                                )
                                gray_last = cv2.cvtColor(
                                    self.previous_frame, cv2.COLOR_BGR2GRAY
                                )
                                
                                # 計算結構相似度
                                score, _ = ssim(
                                    gray_last, gray_current, full=True
                                )
                                
                                # 檢查無變化時間
                                current_time = time.time()
                                time_since_last_change = current_time - self.last_change_time
                                minutes = int(time_since_last_change // 60)
                                seconds = int(time_since_last_change % 60)
                                
                                status_text = f"相似度: {score:.4f} | 無變化時間: {minutes}分{seconds}秒"
                                self.status_var.set(status_text)
                                
                                # 檢查是否超過無變化自動停止時間
                                if (self.auto_stop_enabled and 
                                    self.inactivity_timeout > 0 and 
                                    time_since_last_change > self.inactivity_timeout):
                                    self.log(f"已檢測到 {minutes}分{seconds}秒 無變化，自動停止捕獲")
                                    self.root.after(0, self.stop_capture)
                                    break
                                
                                # 如果幀有明顯變化，保存
                                if score < self.threshold:
                                    self.save_slide(roi_frame)
                                    self.previous_frame = roi_frame
                                    self.last_change_time = current_time  # 更新最後變化時間
                
                    # 在UI中顯示
                    self.root.after(
                        10, lambda: self.display_frame(preview_frame)
                    )
                    
                except Exception as e:
                    self.log(f"捕獲過程中出錯: {str(e)}")
                
                # 等待指定間隔
                time.sleep(self.interval)
    
    def save_slide(self, frame):
        """保存投影片截圖"""
        timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
        filename = os.path.join(
            self.output_folder, 
            f"slide_{timestamp}_{self.slide_count:03d}.png"
        )
        cv2.imwrite(filename, frame)
        self.log(f"保存投影片 #{self.slide_count+1}: {filename}")
        self.slide_count += 1
    
    def pause_capture(self):
        """暫停捕獲"""
        if self.capture_running:
            self.capture_running = False
            self.status_var.set("已暫停")
            self.log("捕獲已暫停")
        else:
            self.capture_running = True
            self.status_var.set("正在捕獲")
            self.log("捕獲已繼續")
    
    def stop_capture(self):
        """停止捕獲"""
        self.capture_running = False
        
        # 更新按鈕狀態
        self.start_btn.config(state=tk.NORMAL)
        self.pause_btn.config(state=tk.DISABLED)
        self.pause_btn.config(text="暫停")
        self.stop_btn.config(state=tk.DISABLED)
        self.go_btn.config(state=tk.NORMAL)
        
        self.status_var.set("已停止")
        self.log(f"捕獲已停止，共獲取 {self.slide_count} 張投影片")
    
    def generate_ppt(self):
        """生成PowerPoint文件"""
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            # 創建演示文稿對象 - 使用16:9比例
            prs = Presentation()
            
            # 設置幻燈片尺寸為16:9 (寬度10英寸, 高度5.625英寸)
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            
            # 檢查投影片目錄中的PNG文件
            slides = sorted([
                f for f in os.listdir(self.output_folder) 
                if f.endswith('.png')
            ])
            
            if not slides:
                messagebox.showinfo("提示", "沒有可用的投影片圖像")
                return
            
            # 添加投影片
            for i, slide_file in enumerate(slides):
                # 添加空白投影片
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白佈局
                
                # 加載投影片圖像
                img_path = os.path.join(self.output_folder, slide_file)
                
                # 獲取圖像尺寸
                img = Image.open(img_path)
                width, height = img.size
                
                # 計算投影片中的位置和大小
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # 保持寬高比例
                ratio = min(slide_width / width, slide_height / height)
                
                # 添加圖片
                left = (slide_width - width * ratio) / 2
                top = (slide_height - height * ratio) / 2
                slide.shapes.add_picture(
                    img_path, left, top, 
                    width=width * ratio, height=height * ratio
                )
            
            # 保存為PowerPoint
            output_path = self.output_file
            prs.save(output_path)
            
            self.log(f"已生成16:9格式PowerPoint文件: {output_path}")
            messagebox.showinfo("成功", f"已生成16:9格式PowerPoint文件: {output_path}")
            
        except Exception as e:
            self.log(f"生成PowerPoint時出錯: {str(e)}")
            messagebox.showerror("錯誤", f"生成PowerPoint時出錯: {str(e)}")
    
    def toggle_auto_stop(self):
        """啟用或禁用自動停止功能"""
        self.auto_stop_enabled = self.auto_stop_var.get()
        self.timeout_scale.config(state=tk.NORMAL if self.auto_stop_enabled else tk.DISABLED)
        self.log(f"自動停止功能已{'啟用' if self.auto_stop_enabled else '禁用'}")
    
    def run(self):
        """運行主循環"""
        self.root.mainloop()
    
    def on_closing(self):
        """窗口關閉事件處理"""
        if self.capture_running:
            if messagebox.askokcancel("退出", "捕獲正在進行中，確定要退出嗎？"):
                self.capture_running = False
                if self.browser:
                    self.browser.quit()
                self.root.destroy()
        else:
            if self.browser:
                self.browser.quit()
            self.root.destroy()

    def create_entry_context_menu(self, entry):
        """為輸入框創建右鍵菜單"""
        # 創建右鍵菜單
        menu = tk.Menu(self.root, tearoff=0)
        menu.add_command(label="剪下", command=lambda: self.entry_cut(entry))
        menu.add_command(label="複製", command=lambda: self.entry_copy(entry))
        menu.add_command(label="貼上", command=lambda: self.entry_paste(entry))
        menu.add_separator()
        menu.add_command(label="全選", command=lambda: self.entry_select_all(entry))
        
        # 綁定右鍵事件
        if sys.platform.startswith('darwin'):  # macOS
            entry.bind("<Button-2>", lambda e: self.show_context_menu(e, menu))
        else:  # Windows/Linux
            entry.bind("<Button-3>", lambda e: self.show_context_menu(e, menu))
    
    def show_context_menu(self, event, menu):
        """顯示右鍵菜單"""
        try:
            menu.tk_popup(event.x_root, event.y_root)
        finally:
            menu.grab_release()
    
    def entry_cut(self, entry):
        """剪下選擇的文字"""
        entry.event_generate("<<Cut>>")
    
    def entry_copy(self, entry):
        """複製選擇的文字"""
        entry.event_generate("<<Copy>>")
    
    def entry_paste(self, entry):
        """貼上文字"""
        entry.event_generate("<<Paste>>")
    
    def entry_select_all(self, entry):
        """全選輸入框內容"""
        entry.select_range(0, tk.END)

    def setup_slide_make_ui(self):
        """設置 Slide 製作標籤頁 UI"""
        frame = self.slide_make_tab
        
        # 頂部說明文字
        info_label = tk.Label(
            frame, 
            text="此功能可從截圖文件夾直接生成 PowerPoint 投影片",
            font=("Arial", 12)
        )
        info_label.pack(pady=20)
        
        # 截圖文件夾選擇區域
        folder_frame = tk.Frame(frame)
        folder_frame.pack(fill=tk.X, pady=10, padx=20)
        
        tk.Label(folder_frame, text="截圖文件夾:").pack(side=tk.LEFT, padx=10)
        self.screenshots_folder_entry = tk.Entry(folder_frame, width=50)
        self.screenshots_folder_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        # 預設填入當前捕獲的輸出文件夾
        self.screenshots_folder_entry.insert(0, self.output_folder)
        
        self.screenshots_browse_btn = tk.Button(
            folder_frame, text="瀏覽...", 
            command=self.browse_screenshots_folder
        )
        self.screenshots_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # PPT 輸出檔案選擇區域
        output_frame = tk.Frame(frame)
        output_frame.pack(fill=tk.X, pady=10, padx=20)
        
        tk.Label(output_frame, text="PowerPoint 輸出:").pack(side=tk.LEFT, padx=10)
        self.ppt_output_entry = tk.Entry(output_frame, width=50)
        self.ppt_output_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        # 預設填入當前的輸出文件
        self.ppt_output_entry.insert(0, os.path.splitext(self.output_file)[0] + ".pptx")
        
        self.ppt_browse_btn = tk.Button(
            output_frame, text="瀏覽...", 
            command=self.browse_ppt_output
        )
        self.ppt_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # 標題輸入區域
        title_frame = tk.Frame(frame)
        title_frame.pack(fill=tk.X, pady=10, padx=20)
        
        tk.Label(title_frame, text="標題:").pack(side=tk.LEFT, padx=10)
        self.ppt_title_entry = tk.Entry(title_frame, width=50)
        self.ppt_title_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        self.ppt_title_entry.insert(0, "視頻捕獲的投影片")
        
        # 圖片排序選項
        sort_frame = tk.Frame(frame)
        sort_frame.pack(fill=tk.X, pady=10, padx=20)
        
        tk.Label(sort_frame, text="圖片排序方式:").pack(side=tk.LEFT, padx=10)
        
        self.name_sort_radio = tk.Radiobutton(
            sort_frame, text="按文件名排序", 
            variable=self.sort_option_var, value="name"
        )
        self.name_sort_radio.pack(side=tk.LEFT, padx=5)
        
        self.time_sort_radio = tk.Radiobutton(
            sort_frame, text="按修改時間排序", 
            variable=self.sort_option_var, value="time"
        )
        self.time_sort_radio.pack(side=tk.LEFT, padx=5)
        
        # 投影片比例選項
        ratio_frame = tk.Frame(frame)
        ratio_frame.pack(fill=tk.X, pady=10, padx=20)
        
        tk.Label(ratio_frame, text="投影片比例:").pack(side=tk.LEFT, padx=10)
        
        self.widescreen_radio = tk.Radiobutton(
            ratio_frame, text="寬屏 (16:9)", 
            variable=self.slide_ratio_var, value="16:9"
        )
        self.widescreen_radio.pack(side=tk.LEFT, padx=5)
        
        self.standard_radio = tk.Radiobutton(
            ratio_frame, text="標準 (4:3)", 
            variable=self.slide_ratio_var, value="4:3"
        )
        self.standard_radio.pack(side=tk.LEFT, padx=5)
        
        # 圖片預覽區域
        preview_frame = tk.LabelFrame(frame, text="圖片預覽")
        preview_frame.pack(fill=tk.X, pady=10, padx=20)
        
        self.preview_canvas = tk.Canvas(preview_frame, height=150, bg="white")
        self.preview_canvas.pack(fill=tk.X, padx=5, pady=5)
        
        # 狀態顯示區域
        status_frame = tk.Frame(frame)
        status_frame.pack(fill=tk.X, pady=5, padx=20)
        
        self.slide_make_status_label = tk.Label(
            status_frame, textvariable=self.slide_make_status_var,
            bd=1, relief=tk.SUNKEN, anchor=tk.W
        )
        self.slide_make_status_label.pack(fill=tk.X, padx=5, pady=5)
        
        # 生成 PPT 按鈕
        self.make_ppt_btn = tk.Button(
            frame, text="生成 PowerPoint", 
            command=self.make_ppt_from_screenshots,
            bg="#4CAF50", fg="white", height=2, width=20
        )
        self.make_ppt_btn.pack(pady=20)
    
    def browse_screenshots_folder(self):
        """瀏覽並選擇截圖文件夾"""
        folder_path = filedialog.askdirectory(
            initialdir=self.output_folder,
            title="選擇包含截圖的文件夾"
        )
        if folder_path:
            self.screenshots_folder_entry.delete(0, tk.END)
            self.screenshots_folder_entry.insert(0, folder_path)
            
            # 更新輸出 PPT 名稱
            folder_name = os.path.basename(folder_path)
            ppt_path = os.path.join(
                os.path.dirname(folder_path),
                f"{folder_name}.pptx"
            )
            
            self.ppt_output_entry.delete(0, tk.END)
            self.ppt_output_entry.insert(0, ppt_path)
            
            # 更新預覽
            self.update_preview()
    
    def browse_ppt_output(self):
        """選擇 PPT 輸出文件"""
        file_path = filedialog.asksaveasfilename(
            initialdir=os.path.dirname(self.output_folder),
            title="保存 PowerPoint 文件",
            filetypes=[("PowerPoint 文件", "*.pptx")]
        )
        if file_path:
            if not file_path.lower().endswith(".pptx"):
                file_path += ".pptx"
                
            self.ppt_output_entry.delete(0, tk.END)
            self.ppt_output_entry.insert(0, file_path)
    
    def update_preview(self):
        """更新圖片預覽"""
        folder_path = self.screenshots_folder_entry.get()
        if not folder_path or not os.path.isdir(folder_path):
            return
            
        # 清除現有預覽
        self.preview_canvas.delete("all")
        
        # 獲取所有圖片文件
        image_files = self.get_image_files_sorted()
        
        if not image_files:
            self.preview_canvas.create_text(
                self.preview_canvas.winfo_width() // 2, 
                self.preview_canvas.winfo_height() // 2,
                text="找不到圖片文件",
                font=("Arial", 14)
            )
            return
            
        # 只顯示前 5 張
        preview_count = min(5, len(image_files))
        preview_images = []
        
        # 計算每張預覽圖的大小和位置
        canvas_width = self.preview_canvas.winfo_width()
        if canvas_width < 100:  # 如果畫布尚未渲染，使用預設寬度
            canvas_width = 500
            
        preview_width = (canvas_width - 10 * (preview_count + 1)) // preview_count
        preview_height = 130
        
        # 載入並顯示預覽圖
        for i in range(preview_count):
            try:
                # 打開圖片並調整大小
                img = Image.open(image_files[i])
                img.thumbnail((preview_width, preview_height), Image.LANCZOS)
                
                # 轉換為 PhotoImage
                photo = ImageTk.PhotoImage(img)
                preview_images.append(photo)  # 保持引用
                
                # 計算位置
                x = 10 + i * (preview_width + 10)
                
                # 在畫布上顯示圖片
                self.preview_canvas.create_image(x, 10, image=photo, anchor=tk.NW)
                
                # 顯示文件名
                filename = os.path.basename(image_files[i])
                if len(filename) > 15:
                    filename = filename[:12] + "..."
                
                self.preview_canvas.create_text(
                    x + preview_width // 2, 
                    preview_height + 20,
                    text=filename,
                    font=("Arial", 8)
                )
                
            except Exception as e:
                print(f"無法載入預覽圖: {e}")
        
        # 保持圖片引用，避免垃圾回收
        self.preview_canvas.preview_images = preview_images
    
    def get_image_files_sorted(self):
        """獲取排序後的圖片文件列表"""
        folder_path = self.screenshots_folder_entry.get()
        if not folder_path or not os.path.isdir(folder_path):
            return []
            
        # 獲取所有圖片文件
        image_files = []
        for filename in os.listdir(folder_path):
            if filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                image_files.append(os.path.join(folder_path, filename))
                
        # 根據選擇的排序方式進行排序
        sort_option = self.sort_option_var.get()
        
        if sort_option == "name":
            # 按文件名排序
            image_files.sort()
        elif sort_option == "time":
            # 按修改時間排序
            image_files.sort(key=os.path.getmtime)
            
        return image_files
    
    def make_ppt_from_screenshots(self):
        """從截圖生成 PowerPoint"""
        folder_path = self.screenshots_folder_entry.get()
        output_path = self.ppt_output_entry.get()
        title = self.ppt_title_entry.get()
        slide_ratio = self.slide_ratio_var.get()
        
        if not folder_path:
            messagebox.showwarning("警告", "請選擇截圖文件夾")
            return
            
        if not os.path.isdir(folder_path):
            messagebox.showerror("錯誤", f"文件夾不存在: {folder_path}")
            return
            
        # 獲取排序後的圖片文件
        image_files = self.get_image_files_sorted()
        
        if not image_files:
            messagebox.showwarning("警告", "選定的文件夾中未找到圖片")
            return
            
        # 確保輸出路徑有效
        if not output_path or output_path.strip() == "" or output_path.strip() == ".pptx":
            # 使用截圖文件夾名稱作為檔案名
            folder_name = os.path.basename(folder_path)
            if not folder_name or folder_name.strip() == "":
                folder_name = "投影片"  # 默認檔案名
                
            output_path = os.path.join(
                os.path.dirname(folder_path),
                f"{folder_name}.pptx"
            )
            self.ppt_output_entry.delete(0, tk.END)
            self.ppt_output_entry.insert(0, output_path)
        elif not output_path.lower().endswith(".pptx"):
            # 確保有正確的副檔名
            output_path += ".pptx"
            self.ppt_output_entry.delete(0, tk.END)
            self.ppt_output_entry.insert(0, output_path)
            
        # 開始生成（在背景線程中運行）
        self.slide_make_status_var.set("正在生成 PowerPoint...")
        self.make_ppt_btn.config(state=tk.DISABLED)
        
        def generate_thread():
            success, result = generate_ppt_from_images(
                folder_path, output_path, title, image_files, slide_ratio
            )
            
            # 在主線程中更新 UI
            self.root.after(
                0, lambda: self.ppt_generation_completed(success, result)
            )
        
        threading.Thread(target=generate_thread).start()
    
    def ppt_generation_completed(self, success, result):
        """PowerPoint 生成完成後的處理"""
        self.make_ppt_btn.config(state=tk.NORMAL)
        
        if success:
            self.slide_make_status_var.set("PowerPoint 生成成功")
            messagebox.showinfo("成功", f"PowerPoint 已成功生成: {result}")
        else:
            self.slide_make_status_var.set(f"PowerPoint 生成失敗: {result}")
            messagebox.showerror("錯誤", f"PowerPoint 生成失敗: {result}")

def main():
    parser = argparse.ArgumentParser(description="Chrome瀏覽器投影片捕獲工具")
    parser.add_argument("--url", help="要打開的網址")
    parser.add_argument(
        "--threshold", type=float, default=0.95, help="相似度閾值 (0.5-1.0)"
    )
    parser.add_argument(
        "--interval", type=float, default=0.5, help="檢測間隔(秒)"
    )
    args = parser.parse_args()
    
    app = ChromeCapture()
    
    # 如果命令行指定了URL，自動打開
    if args.url:
        app.url_entry.insert(0, args.url)
        app.launch_browser()
    
    # 設置命令行參數
    app.threshold_scale.set(args.threshold)
    app.interval_scale.set(args.interval)
    
    app.run()

if __name__ == "__main__":
    main() 