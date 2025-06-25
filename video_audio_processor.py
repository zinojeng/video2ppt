#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
視頻和音頻處理工具

此程序可以處理以下功能：
1. 從 .mp4 文件中提取音頻
2. 捕獲視頻中的幻燈片，並生成 PowerPoint 或 Markdown 文件
"""

import os
import sys
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, simpledialog
import subprocess
import traceback
import threading
from PIL import Image, ImageTk
import importlib
import time
import numpy as np

# 需要時會動態導入的模塊:
# - moviepy (extract_audio_from_video)
# - cv2, numpy, skimage.metrics (capture_slides_from_video)
# - pptx (generate_ppt_from_images)
# - markitdown_helper (custom module for image to markdown conversion)
# - importlib.util (for checking installed modules)
# - selenium, pyautogui, webdriver-manager (web_capture_slides)

def check_dependencies():
    """檢查必要依賴是否已安裝"""
    import importlib.util
    
    required_packages = [
        "opencv-python", "numpy", "pillow", "moviepy",
        "python-pptx", "scikit-image"
    ]
    
    # 分開檢查 markitdown 和 web 捕獲相關套件，因為它們是可選的
    optional_packages = ["markitdown"]
    web_packages = ["selenium", "pyautogui", "webdriver-manager"]
    
    missing_packages = []
    missing_optional = []
    missing_web = []
    
    for package in required_packages:
        try:
            # 處理特殊套件名稱
            import_name = package.replace("-", "_")
            # 處理 opencv-python 特例 (支援 opencv-python 或 opencv-python-headless)
            if package == "opencv-python":
                import_name = "cv2"
                try:
                    # 嘗試導入 cv2 並測試基本功能
                    import cv2
                    if not hasattr(cv2, 'VideoCapture'):
                        missing_packages.append(package)
                    # 如果成功導入，跳過後續檢查
                    continue
                except ImportError:
                    missing_packages.append(package)
                    continue
            # 處理 python-pptx 特例
            elif package == "python-pptx":
                import_name = "pptx"
            # 處理 pillow 特例
            elif package == "pillow":
                import_name = "PIL"
            # 處理 scikit-image 特例
            elif package == "scikit-image":
                import_name = "skimage"
                
            spec = importlib.util.find_spec(import_name)
            if spec is None:
                missing_packages.append(package)
        except ImportError:
            missing_packages.append(package)
        except Exception as e:
            # 新增：OpenCV 導入失敗的特別處理
            if package == "opencv-python":
                print(f"OpenCV 導入失敗: {str(e)}")
                missing_packages.append(package)
    
    for package in optional_packages:
        try:
            spec = importlib.util.find_spec(package)
            if spec is None:
                missing_optional.append(package)
        except ImportError:
            missing_optional.append(package)
    
    for package in web_packages:
        try:
            spec = importlib.util.find_spec(package)
            if spec is None:
                missing_web.append(package)
        except ImportError:
            missing_web.append(package)
    
    if missing_optional:
        print(f"注意: 未安裝可選套件: {', '.join(missing_optional)}")
        print("如果您想使用 MarkItDown 處理功能，請安裝:")
        print("pip install markitdown>=0.1.1")
    
    if missing_web:
        print(f"注意: 未安裝網頁捕獲套件: {', '.join(missing_web)}")
        print("如果您想使用網頁捕獲功能，請安裝:")
        print(f"pip install {' '.join(missing_web)}")
    
    return missing_packages


def install_dependencies(packages):
    """安裝缺失的依賴"""
    print(f"正在安裝必要依賴: {', '.join(packages)}")
    
    try:
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install"] + packages
        )
        return True
    except subprocess.CalledProcessError:
        return False


def extract_audio_from_video(video_path, output_path=None, format="mp3"):
    """
    從視頻文件中提取音頻
    
    參數:
        video_path: 視頻文件路徑
        output_path: 輸出音頻文件路徑，默認為原視頻名稱加上音頻格式後綴
        format: 輸出音頻格式，默認為 mp3
        
    返回:
        success: 是否成功
        output_file: 輸出文件路徑或錯誤信息
    """
    try:
        from moviepy import VideoFileClip
        
        # 如果未指定輸出路徑，使用默認路徑
        if not output_path:
            base_name = os.path.splitext(video_path)[0]
            output_path = f"{base_name}.{format}"
            
        # 打開視頻文件
        video_clip = VideoFileClip(video_path)
        
        # 提取音頻
        audio_clip = video_clip.audio
        
        # 寫入音頻文件
        audio_clip.write_audiofile(output_path)
        
        # 釋放資源
        audio_clip.close()
        video_clip.close()
        
        return True, output_path
        
    except Exception as e:
        error_msg = str(e)
        print(f"提取音頻時出錯: {error_msg}")
        return False, error_msg


def capture_slides_from_video(video_path, output_folder=None, 
                             similarity_threshold=0.7, crop_region=None, 
                             stability_threshold=0.95, sample_interval=0.5,
                             playback_speed=1.0, stop_flag=False, use_fast_mode=True):
    """
    從視頻文件中捕獲幻燈片
    
    參數:
        video_path: 視頻文件路徑
        output_folder: 輸出幻燈片圖片的文件夾
        similarity_threshold: 圖像相似度閾值
        crop_region: 裁剪區域 (x, y, width, height) 或 None
        stability_threshold: 畫面穩定性閾值
        sample_interval: 取樣間隔（秒）
        playback_speed: 播放速度倍數 (1.0=正常速度, 2.0=2倍速, 0.5=0.5倍速)
        stop_flag: 停止捕獲標誌
        use_fast_mode: 是否使用快速相似度計算模式 (True=快速模式, False=精確模式)
    """
    try:
        import cv2
        import numpy as np
        from skimage.metrics import structural_similarity as ssim
        
        # 快速相似度計算函數（優化版本）
        def fast_similarity(img1, img2, use_fast=True):
            """
            快速相似度計算，支援兩種模式：
            - use_fast=True: 使用直方圖+縮小圖像的快速比較
            - use_fast=False: 使用原始完整 SSIM 計算
            """
            if not use_fast:
                # 原始完整 SSIM 計算（保留原功能）
                return ssim(img1, img2)
            
            # 快速模式：先用直方圖篩選
            hist1 = cv2.calcHist([img1], [0], None, [256], [0, 256])
            hist2 = cv2.calcHist([img2], [0], None, [256], [0, 256])
            hist_corr = cv2.compareHist(hist1, hist2, cv2.HISTCMP_CORREL)
            
            # 如果直方圖差異太大，直接返回低相似度
            if hist_corr < 0.7:
                return hist_corr * 0.8  # 調整到 SSIM 相似的範圍
            
            # 縮小圖像到 1/4 尺寸進行 SSIM 計算
            h, w = img1.shape
            small_img1 = cv2.resize(img1, (w//4, h//4))
            small_img2 = cv2.resize(img2, (w//4, h//4))
            
            try:
                small_ssim = ssim(small_img1, small_img2)
                # 對於高相似度的情況，可選擇性地用完整圖像驗證
                if small_ssim > 0.9:
                    return ssim(img1, img2)  # 高相似度時使用完整計算確保準確性
                return small_ssim
            except:
                return hist_corr * 0.8
        
        if not output_folder:
            video_name = os.path.splitext(os.path.basename(video_path))[0]
            output_folder = f"video_slides_{video_name}"
        
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)
            
        cap = cv2.VideoCapture(video_path)
        if not cap.isOpened():
            return False, "無法打開視頻文件"
            
        fps = cap.get(cv2.CAP_PROP_FPS)
        frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        duration = frame_count / fps
        
        # 獲取第一幀用於區域選擇
        ret, first_frame = cap.read()
        if not ret:
            cap.release()
            return False, "無法讀取視頻第一幀"
        
        # 確保裁剪區域有效
        frame_height, frame_width = first_frame.shape[:2]
        if crop_region is None:
            crop_region = (0, 0, frame_width, frame_height)
        else:
            # 確保裁剪區域不超出畫面範圍
            x, y, w, h = crop_region
            x = max(0, min(x, frame_width - 1))
            y = max(0, min(y, frame_height - 1))
            w = min(w, frame_width - x)
            h = min(h, frame_height - y)
            crop_region = (x, y, w, h)
        
        # 重置視頻到開始位置
        cap.set(cv2.CAP_PROP_POS_FRAMES, 0)
        
        frame_idx = 0
        prev_frame = None
        saved_count = 0
        last_saved_time = -1000
        stable_frames = 0
        min_stable_frames = max(1, int(fps * 0.5))  # 增加到 0.5 秒

        # 記錄最後保存的幻燈片內容
        last_saved_slide = None
        
        # 新增：記錄所有已保存的投影片，用於全局比較
        saved_slides = []
        
        # 新增：最小時間間隔（秒），根據播放速度調整
        min_time_between_slides = 2.0 / playback_speed  # 播放速度越快，時間間隔相對越短

        # 使用傳入的取樣間隔，結合播放速度
        effective_interval = sample_interval / playback_speed
        frame_step = max(1, int(fps * sample_interval * playback_speed))
        
        print(f"開始捕獲幻燈片，參數: 相似度={similarity_threshold:.2f}, "
              f"穩定性={stability_threshold:.2f}, 間隔={sample_interval:.2f}秒, "
              f"播放速度={playback_speed:.1f}x")

        while True:
            # 檢查停止標誌
            if stop_flag:
                print("捕獲過程已被用戶停止")
                cap.release()
                return False, "捕獲過程已被用戶停止"
                
            cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
            ret, frame = cap.read()
            if not ret:
                break
                
            current_time = frame_idx / fps
            
            # 裁剪畫面
            x, y, w, h = crop_region
            cropped_frame = frame[y:y+h, x:x+w]
            
            # 轉換為灰度圖像用於相似度比較
            gray_frame = cv2.cvtColor(cropped_frame, cv2.COLOR_BGR2GRAY)
            
            # 畫面穩定性檢測
            is_stable = False
            if prev_frame is not None:
                try:
                    similarity = ssim(prev_frame, gray_frame)
                    if similarity > stability_threshold:
                        stable_frames += 1
                    else:
                        stable_frames = 0
                    
                    # 只有當畫面穩定足夠長時間才考慮保存
                    is_stable = stable_frames >= min_stable_frames
                except ValueError:
                    # 處理圖像尺寸不一致的情況
                    is_stable = False
                    stable_frames = 0
            else:
                # 第一幀總是視為穩定
                is_stable = True
            
            save_condition = False
            
            # 情況1: 第一幀總是保存
            if prev_frame is None:
                save_condition = True
            
            # 情況2: 畫面穩定且滿足條件
            elif is_stable:
                # 檢查時間間隔
                time_since_last_save = current_time - last_saved_time
                if time_since_last_save < min_time_between_slides:
                    save_condition = False
                else:
                    # 與前一幀比較（根據模式選擇計算方法）
                    current_similarity = fast_similarity(prev_frame, gray_frame, use_fast_mode)
                    
                    # 與最後保存的幻燈片比較
                    if last_saved_slide is not None:
                        saved_similarity = fast_similarity(last_saved_slide, gray_frame, use_fast_mode)
                    else:
                        saved_similarity = 0
                    
                    # 新增：檢測內容增加（處理分段顯示的投影片）
                    content_increased = False
                    if last_saved_slide is not None:
                        # 計算非零像素的比例來判斷內容是否增加
                        current_content = np.count_nonzero(gray_frame > 20)
                        last_content = np.count_nonzero(last_saved_slide > 20)
                        content_ratio = current_content / (last_content + 1)
                        
                        # 如果內容增加超過 20%，視為新內容
                        if content_ratio > 1.2:
                            content_increased = True
                            print(f"檢測到內容增加: {content_ratio:.2f}倍")
                    
                    # 新增：與所有已保存的投影片比較
                    is_duplicate = False
                    should_replace = -1  # 記錄應該替換的投影片索引
                    
                    # 限制全局比較數量以提升性能
                    max_compare_slides = min(20, len(saved_slides))  # 最多比較最近20張
                    recent_slides = saved_slides[-max_compare_slides:] if saved_slides else []
                    
                    for i, saved_slide in enumerate(recent_slides):
                        try:
                            # 根據模式選擇相似度計算方法進行全局比較
                            global_similarity = fast_similarity(saved_slide, gray_frame, use_fast_mode)
                            
                            # 如果非常相似但內容增加，應該替換舊的
                            if global_similarity > 0.85 and content_increased:
                                # 檢查當前幀是否包含更多內容
                                saved_content = np.count_nonzero(saved_slide > 20)
                                current_content = np.count_nonzero(gray_frame > 20)
                                if current_content > saved_content * 1.2:
                                    # 計算在原始 saved_slides 中的實際索引
                                    actual_index = len(saved_slides) - max_compare_slides + i
                                    should_replace = actual_index
                                    print(f"將替換投影片 #{actual_index}（內容更完整）")
                                    break
                            
                            # 完全相同的投影片
                            elif global_similarity > 0.95:
                                is_duplicate = True
                                break
                        except Exception:
                            continue
                    
                    # 滿足條件才保存：
                    # 1. 不是重複的投影片
                    # 2. 與前一幀有顯著差異 或 與最後保存的有顯著差異
                    # 3. 時間間隔足夠長
                    # 4. 或者是內容增加的情況
                    if not is_duplicate:
                        save_condition = (
                            current_similarity < similarity_threshold or
                            saved_similarity < similarity_threshold or
                            content_increased or
                            should_replace >= 0
                        )
            
            if save_condition:
                # 如果是替換情況，使用相同的檔名
                if 'should_replace' in locals() and should_replace >= 0:
                    output_path = os.path.join(
                        output_folder, f"slide_{should_replace:03d}.png"
                    )
                    saved_slides[should_replace] = gray_frame.copy()
                    print(f"替換投影片 #{should_replace} 在 {current_time:.2f}秒")
                else:
                    output_path = os.path.join(
                        output_folder, f"slide_{saved_count:03d}.png"
                    )
                    saved_count += 1
                    # 保存到全局列表（限制數量以節省記憶體）
                    if len(saved_slides) > 50:
                        saved_slides.pop(0)
                    saved_slides.append(gray_frame.copy())
                    print(f"保存幻燈片 #{saved_count-1} 在 {current_time:.2f}秒")
                
                cv2.imwrite(output_path, cropped_frame)
                last_saved_time = current_time
                last_saved_slide = gray_frame.copy()
                stable_frames = 0  # 重置穩定性計數器
            
            prev_frame = gray_frame
            frame_idx += frame_step  # 使用設定的取樣間隔
                
            if frame_idx >= frame_count:
                break
                
        cap.release()
        print(f"捕獲完成，共保存 {saved_count} 張幻燈片")
        return True, {
            "output_folder": output_folder,
            "slide_count": saved_count,
            "video_duration": duration
        }
        
    except ImportError as e:
        if "cv2" in str(e):
            error_msg = "未安裝 OpenCV，請執行: pip install opencv-python"
        elif "skimage" in str(e):
            error_msg = "未安裝 scikit-image，請執行: pip install scikit-image"
        else:
            error_msg = str(e)
        return False, error_msg
    except Exception as e:
        error_msg = str(e)
        traceback.print_exc()
        return False, error_msg


def generate_ppt_from_images(image_folder, output_file=None, title="視頻捕獲的幻燈片", 
                             image_files=None, slide_ratio="16:9"):
    """
    將圖片文件夾轉換為 PowerPoint 文件
    
    參數:
        image_folder: 包含幻燈片圖片的文件夾
        output_file: 輸出的 PowerPoint 文件路徑，默認為文件夾名+.pptx
        title: 演示文稿標題
        image_files: 已排序的圖片文件列表，若為 None 則自動從文件夾獲取
        slide_ratio: 幻燈片比例，支持 "16:9" 或 "4:3"
        
    返回:
        success: 是否成功
        output_file: 輸出文件路徑或錯誤信息
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        # 如果未指定輸出文件，使用默認文件
        if not output_file:
            output_file = os.path.join(
                os.path.dirname(image_folder), 
                f"{os.path.basename(image_folder)}.pptx"
            )
        
        # 創建新的 PowerPoint 演示文稿，根據選擇的比例
        prs = Presentation()
        
        # 設置幻燈片大小 (預設單位是 EMU, 1 英寸 = 914400 EMU)
        if slide_ratio == "16:9":
            prs.slide_width = 9144000  # 10 英寸 (16 單位寬)
            prs.slide_height = 5143500  # 5.625 英寸 (9 單位高)
        elif slide_ratio == "4:3":
            prs.slide_width = 9144000  # 10 英寸 (4 單位寬)
            prs.slide_height = 6858000  # 7.5 英寸 (3 單位高)
        
        # 獲取幻燈片尺寸
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # 如果未提供圖片文件列表，則從文件夾獲取
        if image_files is None:
            image_files = []
            for filename in sorted(os.listdir(image_folder)):
                if filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                    image_files.append(os.path.join(image_folder, filename))
                
        if not image_files:
            return False, "未找到圖片文件"
            
        # 添加標題幻燈片
        title_slide_layout = prs.slide_layouts[0]  # 標題布局
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = title
        if hasattr(title_slide.placeholders, 'subtitle'):
            subtitle = f"包含 {len(image_files)} 張幻燈片 ({slide_ratio})"
            title_slide.placeholders[1].text = subtitle
            
        # 為每張圖片創建一個幻燈片
        blank_slide_layout = prs.slide_layouts[6]  # 空白布局
        
        for img_path in image_files:
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加圖片
            left = top = Inches(0)
            slide.shapes.add_picture(
                img_path, left, top, width=slide_width, height=slide_height
            )
            
        # 保存演示文稿
        prs.save(output_file)
        
        return True, output_file
        
    except Exception as e:
        error_msg = str(e)
        print(f"生成 PowerPoint 時出錯: {error_msg}")
        traceback.print_exc()
        return False, error_msg


def generate_markdown_from_images(image_folder, output_file=None, title="視頻捕獲的幻燈片", use_markitdown=True, api_key=None):
    """
    將圖片文件夾轉換為 Markdown 文件
    
    參數:
        image_folder: 包含幻燈片圖片的文件夾
        output_file: 輸出的 Markdown 文件路徑，默認為文件夾名+.md
        title: Markdown 文件標題
        use_markitdown: 是否使用 MarkItDown 庫進行圖片文本提取
        api_key: 如果使用 MarkItDown 並需要 LLM 支持，提供 API Key
        
    返回:
        success: 是否成功
        output_file: 輸出文件路徑或錯誤信息
    """
    try:
        # 如果未指定輸出文件，使用默認文件
        if not output_file:
            output_file = os.path.join(
                os.path.dirname(image_folder), 
                f"{os.path.basename(image_folder)}.md"
            )
            
        # 獲取所有圖片文件
        image_files = []
        for filename in sorted(os.listdir(image_folder)):
            if filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                image_files.append(os.path.join(image_folder, filename))
                
        if not image_files:
            return False, "未找到圖片文件"
        
        # 嘗試使用我們的自定義 markitdown_helper
        try:
            import markitdown_helper
            
            print(f"使用自定義 markitdown_helper 處理 {len(image_files)} 張圖片...")
            success, result, info = markitdown_helper.convert_images_to_markdown(
                image_paths=image_files,
                output_file=output_file,
                title=title,
                use_llm=(api_key is not None),
                api_key=api_key
            )
            
            if success:
                return True, result
            else:
                print(f"使用 markitdown_helper 失敗: {info.get('error', '未知錯誤')}")
                # 繼續使用基本方法
        except ImportError:
            print("找不到 markitdown_helper 模組，繼續嘗試其他方法...")
        except Exception as e:
            print(f"使用 markitdown_helper 時出錯: {e}")
            traceback.print_exc()
            
        # 如果使用 MarkItDown，則調用相應函數
        if use_markitdown:
            try:
                # 檢查 markitdown 是否已安裝
                import importlib.util
                markitdown_spec = importlib.util.find_spec("markitdown")
                
                if markitdown_spec is not None:
                    import markitdown
                    
                    # 使用 markitdown 庫
                    print(f"使用 MarkItDown 庫處理 {len(image_files)} 張圖片...")
                    
                    # 創建 MarkItDown 實例
                    converter = markitdown.MarkItDown()
                    
                    # 添加自訂標題的 Markdown 內容
                    with open(output_file, 'w', encoding='utf-8') as f:
                        f.write(f"# {title}\n\n")
                        
                        for i, img_path in enumerate(image_files):
                            # 獲取相對路徑
                            try:
                                rel_path = os.path.relpath(
                                    img_path,
                                    os.path.dirname(output_file)
                                )
                            except Exception:
                                rel_path = img_path
                                
                            # 添加幻燈片標題和圖片
                            slide_num = i + 1
                            f.write(f"## 幻燈片 {slide_num}\n\n")
                            f.write(f"![幻燈片 {slide_num}]({rel_path})\n\n")
                            
                            try:
                                # 使用 markitdown 分析圖片內容
                                result = converter.convert(img_path)
                                if result and result.text_content:
                                    f.write(f"{result.text_content.strip()}\n\n")
                            except Exception as img_err:
                                print(f"分析圖片 {img_path} 時出錯: {img_err}")
                            
                            f.write("---\n\n")
                    
                    print(f"已成功生成 Markdown 文件: {output_file}")
                    return True, output_file
                else:
                    print("找不到 markitdown 模組，將使用基本方法生成 Markdown")
                    use_markitdown = False
                    
            except ImportError as e:
                print(f"MarkItDown 庫導入錯誤: {e}")
                print("將使用基本方法生成 Markdown")
                use_markitdown = False
            except Exception as e:
                print(f"使用 MarkItDown 時出錯: {e}")
                print("將使用基本方法生成 Markdown")
                traceback.print_exc()
                use_markitdown = False
        
        # 使用基本方法生成 Markdown
        print("使用基本方法生成 Markdown...")
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(f"# {title}\n\n")
            
            for i, img_path in enumerate(image_files):
                # 獲取相對路徑，使 Markdown 中的圖片鏈接更短
                try:
                    rel_path = os.path.relpath(
                        img_path, 
                        os.path.dirname(output_file)
                    )
                except Exception:
                    rel_path = img_path
                    
                # 添加幻燈片標題和圖片
                slide_num = i + 1
                f.write(f"## 幻燈片 {slide_num}\n\n")
                f.write(f"![幻燈片 {slide_num}]({rel_path})\n\n")
                f.write("---\n\n")
        
        return True, output_file
        
    except Exception as e:
        error_msg = str(e)
        print(f"生成 Markdown 時出錯: {error_msg}")
        traceback.print_exc()
        return False, error_msg


class VideoAudioProcessor:
    """視頻和音頻處理應用"""
    
    def __init__(self, root):
        self.root = root
        self.root.title("視頻和音頻處理工具")
        self.root.geometry("900x700")
        
        # 保存使用者的 API Key
        self.saved_api_key = os.environ.get("OPENAI_API_KEY", "")
        
        # 創建頁面框架
        self.setup_ui()
        
        self.crop_region = None  # 新增：儲存裁剪區域
    
    def setup_ui(self):
        """設置使用者介面"""
        # 建立選項卡
        self.notebook = ttk.Notebook(self.root)
        self.notebook.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        # 第一個選項卡：音頻提取
        self.audio_frame = ttk.Frame(self.notebook)
        self.notebook.add(self.audio_frame, text="音頻提取")
        
        # 第二個選項卡：幻燈片捕獲
        self.slide_frame = ttk.Frame(self.notebook)
        self.notebook.add(self.slide_frame, text="幻燈片捕獲")
        
        # 第三個選項卡：幻燈片處理
        self.process_frame = ttk.Frame(self.notebook)
        self.notebook.add(self.process_frame, text="幻燈片處理")
        
        # 第四個選項卡：PPT Slide 製作（從截圖文件夾生成）
        self.web_capture_frame = ttk.Frame(self.notebook)
        self.notebook.add(self.web_capture_frame, text="PPT Slide 製作")
        
        # 設置音頻提取 UI
        self.setup_audio_ui()
        
        # 設置幻燈片捕獲 UI
        self.setup_slide_ui()
        
        # 設置幻燈片處理 UI
        self.setup_process_ui()
        
        # 設置網頁捕獲 UI
        self.setup_web_capture_ui()
    
    def setup_audio_ui(self):
        """設置音頻提取界面"""
        frame = self.audio_frame
        
        # 頂部說明文字
        info_label = tk.Label(
            frame, 
            text="此功能可以從視頻文件中提取音頻軌道",
            font=("Arial", 12)
        )
        info_label.pack(pady=20)
        
        # 視頻文件選擇區域
        video_frame = tk.Frame(frame)
        video_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(video_frame, text="視頻文件:").pack(side=tk.LEFT, padx=10)
        self.video_entry = tk.Entry(video_frame, width=50)
        self.video_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        self.video_browse_btn = tk.Button(
            video_frame, text="瀏覽...", 
            command=lambda: self.browse_file(
                self.video_entry, 
                filetypes=[("視頻文件", "*.mp4 *.avi *.mkv *.mov")]
            )
        )
        self.video_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # 輸出音頻設置區域
        output_frame = tk.Frame(frame)
        output_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(output_frame, text="輸出文件:").pack(side=tk.LEFT, padx=10)
        self.audio_output_entry = tk.Entry(output_frame, width=50)
        self.audio_output_entry.pack(
            side=tk.LEFT, padx=5, fill=tk.X, expand=True
        )
        
        self.audio_browse_btn = tk.Button(
            output_frame, text="瀏覽...", 
            command=lambda: self.save_file(
                self.audio_output_entry, 
                filetypes=[("音頻文件", "*.mp3 *.wav *.aac")]
            )
        )
        self.audio_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # 音頻格式選擇
        format_frame = tk.Frame(frame)
        format_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(format_frame, text="音頻格式:").pack(side=tk.LEFT, padx=10)
        
        self.audio_format_var = tk.StringVar(value="mp3")
        
        tk.Radiobutton(
            format_frame, text="MP3", 
            variable=self.audio_format_var, value="mp3"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            format_frame, text="WAV", 
            variable=self.audio_format_var, value="wav"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            format_frame, text="AAC", 
            variable=self.audio_format_var, value="aac"
        ).pack(side=tk.LEFT, padx=5)
        
        # 狀態顯示
        self.audio_status_var = tk.StringVar(value="準備就緒")
        status_label = tk.Label(
            frame, textvariable=self.audio_status_var,
            font=("Arial", 10), fg="blue"
        )
        status_label.pack(pady=10)
        
        # 進度條
        self.audio_progress = ttk.Progressbar(
            frame, orient="horizontal", length=300, mode="indeterminate"
        )
        self.audio_progress.pack(pady=10)
        
        # 提取按鈕
        self.extract_btn = tk.Button(
            frame, text="提取音頻", command=self.extract_audio,
            bg="#4CAF50", fg="white", height=2, width=20
        )
        self.extract_btn.pack(pady=20)
    
    def browse_file(self, entry_widget, filetypes):
        """瀏覽並選擇文件"""
        file_path = filedialog.askopenfilename(
            initialdir=".",
            title="選擇文件",
            filetypes=filetypes
        )
        if file_path:
            entry_widget.delete(0, tk.END)
            entry_widget.insert(0, file_path)
            
            # 如果是視頻文件，自動生成音頻輸出路徑
            if entry_widget == self.video_entry:
                base_name = os.path.splitext(file_path)[0]
                audio_format = self.audio_format_var.get()
                output_path = f"{base_name}.{audio_format}"
                
                self.audio_output_entry.delete(0, tk.END)
                self.audio_output_entry.insert(0, output_path)
    
    def save_file(self, entry_widget, filetypes):
        """選擇保存文件的路徑"""
        file_path = filedialog.asksaveasfilename(
            initialdir=".",
            title="保存文件",
            filetypes=filetypes
        )
        if file_path:
            entry_widget.delete(0, tk.END)
            entry_widget.insert(0, file_path)
    
    def extract_audio(self):
        """提取音頻按鈕點擊處理"""
        video_path = self.video_entry.get()
        output_path = self.audio_output_entry.get()
        audio_format = self.audio_format_var.get()
        
        if not video_path:
            messagebox.showwarning("警告", "請選擇視頻文件")
            return
            
        # 確保輸出路徑有正確的擴展名
        if output_path and not output_path.lower().endswith(f".{audio_format}"):
            output_path = f"{os.path.splitext(output_path)[0]}.{audio_format}"
            self.audio_output_entry.delete(0, tk.END)
            self.audio_output_entry.insert(0, output_path)
        
        # 開始提取（在背景線程中運行）
        self.audio_status_var.set("正在提取音頻...")
        self.audio_progress.start(10)
        self.extract_btn.config(state=tk.DISABLED)
        
        def extract_thread():
            success, result = extract_audio_from_video(
                video_path, output_path, audio_format
            )
            
            # 在主線程中更新 UI
            self.root.after(0, lambda: self.extraction_completed(success, result))
        
        threading.Thread(target=extract_thread).start()
    
    def extraction_completed(self, success, result):
        """音頻提取完成後的處理"""
        self.audio_progress.stop()
        self.extract_btn.config(state=tk.NORMAL)
        
        if success:
            self.audio_status_var.set(f"提取完成: {result}")
            messagebox.showinfo("成功", f"音頻已成功提取到: {result}")
        else:
            self.audio_status_var.set(f"提取失敗: {result}")
            messagebox.showerror("錯誤", f"音頻提取失敗: {result}")
            
    def setup_slide_ui(self):
        """設置幻燈片捕獲界面"""
        frame = self.slide_frame
        
        # 頂部說明文字
        info_label = tk.Label(
            frame, 
            text="此功能可從視頻文件中捕獲幻燈片，並將其保存為圖片",
            font=("Arial", 12)
        )
        info_label.pack(pady=20)
        
        # 在頂部添加安裝提示
        install_label = tk.Label(
            frame, 
            text="注意: 此功能需要安裝 OpenCV (pip install opencv-python)",
            font=("Arial", 10), fg="red"
        )
        install_label.pack(pady=5)
        
        # 視頻文件選擇區域
        video_frame = tk.Frame(frame)
        video_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(video_frame, text="視頻文件:").pack(side=tk.LEFT, padx=10)
        self.slide_video_entry = tk.Entry(video_frame, width=50)
        self.slide_video_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        self.slide_video_browse_btn = tk.Button(
            video_frame, text="瀏覽...", 
            command=lambda: self.browse_file(
                self.slide_video_entry, 
                filetypes=[("視頻文件", "*.mp4 *.avi *.mkv *.mov")]
            )
        )
        self.slide_video_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # 輸出文件夾設置區域
        output_frame = tk.Frame(frame)
        output_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(output_frame, text="輸出文件夾:").pack(side=tk.LEFT, padx=10)
        self.slide_output_entry = tk.Entry(output_frame, width=50)
        self.slide_output_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        self.slide_browse_btn = tk.Button(
            output_frame, text="瀏覽...", 
            command=self.browse_slide_folder
        )
        self.slide_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # === 新增：參數調整框架 ===
        params_frame = ttk.LabelFrame(frame, text="捕獲參數調整")
        params_frame.pack(fill=tk.X, pady=10, padx=10)
        
        # 相似度閾值設置
        threshold_frame = tk.Frame(params_frame)
        threshold_frame.pack(fill=tk.X, pady=5)
        
        tk.Label(threshold_frame, text="相似度閾值:").pack(side=tk.LEFT, padx=10)
        
        self.threshold_var = tk.DoubleVar(value=0.85)  # 提高預設值
        self.threshold_scale = ttk.Scale(
            threshold_frame, from_=0.5, to=0.95, 
            variable=self.threshold_var, 
            length=200,
            orient=tk.HORIZONTAL
        )
        self.threshold_scale.pack(side=tk.LEFT, padx=5)
        
        # 顯示當前閾值
        self.threshold_label = tk.Label(
            threshold_frame, 
            text=f"{self.threshold_var.get():.2f}"  # 初始值顯示
        )
        self.threshold_label.pack(side=tk.LEFT, padx=5)
        
        # 新增說明文字
        threshold_tip = tk.Label(
            threshold_frame, 
            text="值越低越靈敏，捕獲更多幻燈片（可能包含相似畫面）",
            font=("Arial", 9), fg="#666"
        )
        threshold_tip.pack(side=tk.LEFT, padx=10)
        
        # 更新閾值顯示 (修正綁定)
        self.threshold_var.trace_add("write", self.update_threshold_label)
        
        # 穩定性閾值設置
        stability_frame = tk.Frame(params_frame)
        stability_frame.pack(fill=tk.X, pady=5)
        
        tk.Label(stability_frame, text="穩定性閾值:").pack(side=tk.LEFT, padx=10)
        
        self.stability_var = tk.DoubleVar(value=0.96)  # 提高預設值
        self.stability_scale = ttk.Scale(
            stability_frame, from_=0.85, to=0.99, 
            variable=self.stability_var, 
            length=200,
            orient=tk.HORIZONTAL
        )
        self.stability_scale.pack(side=tk.LEFT, padx=5)
        
        # 顯示當前穩定性閾值
        self.stability_label = tk.Label(
            stability_frame, 
            text=f"{self.stability_var.get():.2f}"  # 初始值顯示
        )
        self.stability_label.pack(side=tk.LEFT, padx=5)
        
        # 更新穩定性閾值顯示 (修正綁定)
        self.stability_var.trace_add("write", self.update_stability_label)
        
        # 說明文字
        stability_tip = tk.Label(
            stability_frame, 
            text="值越高越嚴格，避免殘影（可能錯過快速變化的畫面）",
            font=("Arial", 9), fg="#666"
        )
        stability_tip.pack(side=tk.LEFT, padx=10)
        
        # 取樣間隔設置
        interval_frame = tk.Frame(params_frame)
        interval_frame.pack(fill=tk.X, pady=5)
        
        tk.Label(interval_frame, text="取樣間隔(秒):").pack(side=tk.LEFT, padx=10)
        
        self.interval_var = tk.DoubleVar(value=1.0)  # 增加預設值
        self.interval_scale = ttk.Scale(
            interval_frame, from_=0.1, to=3.0,  # 擴大範圍到 3.0
            variable=self.interval_var, 
            length=200,
            orient=tk.HORIZONTAL
        )
        self.interval_scale.pack(side=tk.LEFT, padx=5)
        
        # 顯示當前取樣間隔
        self.interval_label = tk.Label(
            interval_frame, 
            text=f"{self.interval_var.get():.2f}"  # 初始值顯示
        )
        self.interval_label.pack(side=tk.LEFT, padx=5)
        
        # 更新間隔顯示 (修正綁定)
        self.interval_var.trace_add("write", self.update_interval_label)
        
        # 說明文字
        interval_tip = tk.Label(
            interval_frame, 
            text="值越大效率越高，值越小捕獲更精確",
            font=("Arial", 9), fg="#666"
        )
        interval_tip.pack(side=tk.LEFT, padx=10)
        
        # 播放速度設置
        speed_frame = tk.Frame(params_frame)
        speed_frame.pack(fill=tk.X, pady=5)
        
        tk.Label(speed_frame, text="播放速度:").pack(side=tk.LEFT, padx=10)
        
        self.speed_var = tk.DoubleVar(value=1.0)
        self.speed_scale = ttk.Scale(
            speed_frame, from_=0.5, to=4.0,
            variable=self.speed_var, 
            length=200,
            orient=tk.HORIZONTAL
        )
        self.speed_scale.pack(side=tk.LEFT, padx=5)
        
        # 顯示當前播放速度
        self.speed_label = tk.Label(
            speed_frame, 
            text=f"{self.speed_var.get():.1f}x"
        )
        self.speed_label.pack(side=tk.LEFT, padx=5)
        
        # 更新速度顯示
        self.speed_var.trace_add("write", self.update_speed_label)
        
        # 說明文字
        speed_tip = tk.Label(
            speed_frame, 
            text="值越大截取速度越快（2x = 2倍速，4x = 4倍速）",
            font=("Arial", 9), fg="#666"
        )
        speed_tip.pack(side=tk.LEFT, padx=10)
        
        # 快速模式選項
        fast_mode_frame = tk.Frame(params_frame)
        fast_mode_frame.pack(fill=tk.X, pady=5)
        
        self.fast_mode_var = tk.BooleanVar(value=True)  # 預設啟用快速模式
        fast_mode_check = tk.Checkbutton(
            fast_mode_frame, 
            text="啟用快速模式", 
            variable=self.fast_mode_var,
            font=("Arial", 10)
        )
        fast_mode_check.pack(side=tk.LEFT, padx=10)
        
        fast_mode_tip = tk.Label(
            fast_mode_frame, 
            text="快速模式可提高處理速度3-5倍，保持高精確度",
            font=("Arial", 9), fg="#666"
        )
        fast_mode_tip.pack(side=tk.LEFT, padx=10)
        
        # 參數重置按鈕
        reset_params_btn = tk.Button(
            params_frame, text="重置參數", 
            command=self.reset_slide_params,
            bg="#f0f0f0"
        )
        reset_params_btn.pack(pady=5)
        
        # === 區域選擇按鈕 ===
        region_frame = tk.Frame(frame)
        region_frame.pack(fill=tk.X, pady=10)
        
        self.select_region_btn = tk.Button(
            region_frame, text="選擇幻燈片區域", 
            command=self.select_slide_region,
            state=tk.DISABLED
        )
        self.select_region_btn.pack(side=tk.LEFT, padx=10)
        
        # 區域顯示標籤
        self.region_label = tk.Label(
            region_frame, text="未選擇區域", fg="gray"
        )
        self.region_label.pack(side=tk.LEFT, padx=10)
        
        # 清除區域按鈕
        self.clear_region_btn = tk.Button(
            region_frame, text="清除區域", 
            command=self.clear_slide_region,
            state=tk.DISABLED
        )
        self.clear_region_btn.pack(side=tk.LEFT, padx=5)
        
        # 狀態顯示
        self.slide_status_var = tk.StringVar(value="準備就緒")
        status_label = tk.Label(
            frame, textvariable=self.slide_status_var,
            font=("Arial", 10), fg="blue"
        )
        status_label.pack(pady=10)
        
        # 進度條
        self.slide_progress = ttk.Progressbar(
            frame, orient="horizontal", length=300, mode="indeterminate"
        )
        self.slide_progress.pack(pady=10)
        
        # 捕獲按鈕區域 (修改為按鈕群組)
        button_frame = tk.Frame(frame)
        button_frame.pack(pady=20)
        
        # 捕獲按鈕
        self.capture_btn = tk.Button(
            button_frame, text="捕獲幻燈片", command=self.capture_slides,
            bg="#2196F3", fg="white", height=2, width=15
        )
        self.capture_btn.pack(side=tk.LEFT, padx=10)
        
        # 新增：停止按鈕
        self.stop_btn = tk.Button(
            button_frame, text="停止捕獲", command=self.stop_capture,
            bg="#F44336", fg="white", height=2, width=10,
            state=tk.DISABLED
        )
        self.stop_btn.pack(side=tk.LEFT, padx=10)
        
        # 新增：重置按鈕
        self.reset_btn = tk.Button(
            button_frame, text="重置參數", command=self.reset_slide_params,
            bg="#FFC107", fg="black", height=2, width=10
        )
        self.reset_btn.pack(side=tk.LEFT, padx=10)
        
        # === 新增：預設設定框架 ===
        preset_frame = tk.Frame(params_frame)
        preset_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(preset_frame, text="預設設定:").pack(side=tk.LEFT, padx=10)
        
        self.preset_var = tk.StringVar()
        presets = [
            "標準設定 (推薦)",
            "動畫較多的簡報",
            "靜態畫面為主的簡報",
            "長影片優化",
            "捕獲較少的設定 (靈敏度低)",
            "快速模式 (4倍速)"  # 新增快速模式
        ]
        self.preset_combo = ttk.Combobox(
            preset_frame, textvariable=self.preset_var,
            values=presets, state="readonly", width=25
        )
        self.preset_combo.pack(side=tk.LEFT, padx=5)
        self.preset_combo.bind("<<ComboboxSelected>>", self.apply_preset)
        
        # 添加說明標籤
        self.preset_tip = tk.Label(
            preset_frame, 
            text="選擇預設設定後將自動調整參數",
            font=("Arial", 9), fg="#666"
        )
        self.preset_tip.pack(side=tk.LEFT, padx=10)
        
    def reset_slide_params(self):
        """重置捕獲參數為默認值"""
        self.threshold_var.set(0.85)  # 提高預設值
        self.stability_var.set(0.96)  # 提高預設值
        self.interval_var.set(1.0)    # 增加預設值
        self.speed_var.set(1.0)       # 重置播放速度
        self.fast_mode_var.set(True)  # 重置快速模式為啟用
        messagebox.showinfo("參數重置", "捕獲參數已重置為默認值")
        
    def check_video_selected(self, event=None):
        """檢查是否已選擇視頻文件"""
        video_path = self.slide_video_entry.get()
        if video_path and os.path.isfile(video_path):
            try:
                # 新增：檢查 OpenCV 是否安裝
                import cv2
                self.select_region_btn.config(state=tk.NORMAL)
            except ImportError:
                # 如果未安裝，禁用區域選擇按鈕並顯示提示
                self.select_region_btn.config(state=tk.DISABLED)
                self.region_label.config(
                    text="請先安裝 OpenCV (pip install opencv-python)", 
                    fg="red"
                )
        else:
            self.select_region_btn.config(state=tk.DISABLED)
            self.clear_region_btn.config(state=tk.DISABLED)
            self.region_label.config(text="未選擇區域", fg="gray")
            self.crop_region = None
            
    def select_slide_region(self):
        """讓使用者選擇裁剪區域（新增時間點選擇功能）"""
        video_path = self.slide_video_entry.get()
        if not video_path:
            return
            
        try:
            import cv2
            
            # 彈出對話框讓使用者輸入時間點（秒）
            time_str = simpledialog.askfloat(
                "選擇時間點",
                "請輸入影片時間點（秒）來選取區域:",
                parent=self.root,
                minvalue=0
            )
            
            if time_str is None:  # 使用者取消
                return
                
            # 讀取指定時間點的幀
            cap = cv2.VideoCapture(video_path)
            
            # 設置到指定時間（毫秒）
            cap.set(cv2.CAP_PROP_POS_MSEC, time_str * 1000)
            
            ret, frame = cap.read()
            cap.release()
            
            if not ret:
                messagebox.showerror("錯誤", f"無法讀取影片在 {time_str} 秒的畫面")
                return
                
            # 轉換為 RGB 格式
            frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            
            # 創建區域選擇窗口
            region_selector = tk.Toplevel(self.root)
            region_selector.title(f"選擇幻燈片區域 (時間點: {time_str}秒)")
            region_selector.geometry(f"{frame.shape[1]}x{frame.shape[0]}")
            
            # 顯示圖像
            img = Image.fromarray(frame_rgb)
            photo = ImageTk.PhotoImage(img)
            
            canvas = tk.Canvas(region_selector, width=frame.shape[1], height=frame.shape[0])
            canvas.pack()
            canvas.create_image(0, 0, anchor=tk.NW, image=photo)
            canvas.image = photo  # 保持引用
            
            # 區域選擇變量
            start_x, start_y = None, None
            rect_id = None
            selected_region = None
            
            def on_press(event):
                nonlocal start_x, start_y, rect_id
                start_x, start_y = event.x, event.y
                rect_id = canvas.create_rectangle(
                    start_x, start_y, start_x, start_y, outline="red", width=2
                )
            
            def on_drag(event):
                nonlocal start_x, start_y, rect_id
                if rect_id:
                    canvas.coords(rect_id, start_x, start_y, event.x, event.y)
            
            def on_release(event):
                nonlocal selected_region
                x1, y1 = start_x, start_y
                x2, y2 = event.x, event.y
                
                # 確保左上角和右下角
                x = min(x1, x2)
                y = min(y1, y2)
                width = abs(x2 - x1)
                height = abs(y2 - y1)
                
                selected_region = (x, y, width, height)
                self.crop_region = selected_region
                self.region_label.config(
                    text=f"區域: ({x}, {y}, {width}, {height})",
                    fg="green"
                )
                self.clear_region_btn.config(state=tk.NORMAL)
                region_selector.destroy()
            
            # 綁定事件
            canvas.bind("<ButtonPress-1>", on_press)
            canvas.bind("<B1-Motion>", on_drag)
            canvas.bind("<ButtonRelease-1>", on_release)
            
            # 添加說明文字
            canvas.create_text(
                frame.shape[1] // 2, 20,
                text="拖拽滑鼠選擇幻燈片區域",
                fill="yellow",
                font=("Arial", 16, "bold")
            )
            
        except ImportError:
            messagebox.showerror("錯誤", "請先安裝OpenCV: pip install opencv-python")
        except Exception as e:
            messagebox.showerror("錯誤", f"區域選擇失敗: {str(e)}")
            
    def clear_slide_region(self):
        """清除已選擇的區域"""
        self.crop_region = None
        self.region_label.config(text="未選擇區域", fg="gray")
        self.clear_region_btn.config(state=tk.DISABLED)
            
    def capture_slides(self):
        """捕獲幻燈片按鈕點擊處理"""
        video_path = self.slide_video_entry.get()
        output_folder = self.slide_output_entry.get()
        threshold = self.threshold_var.get()
        stability = self.stability_var.get()
        interval = self.interval_var.get()
        speed = self.speed_var.get()
        
        if not video_path:
            messagebox.showwarning("警告", "請選擇視頻文件")
            return
            
        # 如果未指定輸出文件夾，根據視頻名稱創建
        if not output_folder:
            video_name = os.path.splitext(os.path.basename(video_path))[0]
            output_folder = f"video_slides_{video_name}"
            
            self.slide_output_entry.delete(0, tk.END)
            self.slide_output_entry.insert(0, output_folder)
        
        # 顯示使用的參數
        fast_mode_text = "快速" if self.fast_mode_var.get() else "精確"
        self.slide_status_var.set(
            f"開始捕獲: 相似度={threshold:.2f}, 穩定性={stability:.2f}, 間隔={interval:.2f}秒, 速度={speed:.1f}x, 模式={fast_mode_text}"
        )
        
        # 開始捕獲（在背景線程中運行）
        self.slide_progress.start(10)
        self.capture_btn.config(state=tk.DISABLED)
        self.stop_btn.config(state=tk.NORMAL)
        
        # 設置停止標誌
        self.stop_capture_flag = False
        
        def capture_thread():
            # 調用捕獲函數，傳遞所有參數
            success, result = capture_slides_from_video(
                video_path, 
                output_folder, 
                similarity_threshold=threshold,
                crop_region=self.crop_region,
                stability_threshold=stability,
                sample_interval=interval,
                playback_speed=speed,
                stop_flag=self.stop_capture_flag,  # 傳遞停止標誌
                use_fast_mode=self.fast_mode_var.get()  # 傳遞快速模式設定
            )
            
            # 在主線程中更新 UI
            self.root.after(0, lambda: self.capture_completed(success, result))
        
        self.capture_thread = threading.Thread(target=capture_thread)
        self.capture_thread.daemon = True
        self.capture_thread.start()
    
    def capture_completed(self, success, result):
        """幻燈片捕獲完成後的處理"""
        self.slide_progress.stop()
        self.capture_btn.config(state=tk.NORMAL)
        self.stop_btn.config(state=tk.DISABLED)
        
        if success:
            output_folder = result.get("output_folder", "未知文件夾")
            slide_count = result.get("slide_count", 0)
            
            self.slide_status_var.set(f"捕獲完成: {slide_count} 張幻燈片")
            
            # 同時更新處理標籤頁的文件夾路徑
            self.process_folder_entry.delete(0, tk.END)
            self.process_folder_entry.insert(0, output_folder)
            
            # 詢問是否切換到處理標籤頁
            if messagebox.askyesno(
                "成功", 
                f"已成功捕獲 {slide_count} 張幻燈片到文件夾: {output_folder}\n"
                f"是否立即處理這些幻燈片？"
            ):
                self.notebook.select(2)  # 切換到處理標籤頁
        else:
            self.slide_status_var.set(f"捕獲失敗: {result}")
            messagebox.showerror("錯誤", f"幻燈片捕獲失敗: {result}")

    def setup_process_ui(self):
        """設置幻燈片處理界面"""
        frame = self.process_frame
        
        # 頂部說明文字
        info_label = tk.Label(
            frame, 
            text="此功能可將捕獲的幻燈片轉換為 PowerPoint 或 Markdown",
            font=("Arial", 12)
        )
        info_label.pack(pady=20)
        
        # 幻燈片文件夾選擇區域
        folder_frame = tk.Frame(frame)
        folder_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(folder_frame, text="幻燈片文件夾:").pack(side=tk.LEFT, padx=10)
        self.process_folder_entry = tk.Entry(folder_frame, width=50)
        self.process_folder_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        self.process_browse_btn = tk.Button(
            folder_frame, text="瀏覽...", 
            command=self.browse_process_folder
        )
        self.process_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # 輸出格式選擇
        format_frame = tk.Frame(frame)
        format_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(format_frame, text="輸出格式:").pack(side=tk.LEFT, padx=10)
        
        self.output_format_var = tk.StringVar(value="both")
        
        tk.Radiobutton(
            format_frame, text="僅 PowerPoint", 
            variable=self.output_format_var, value="pptx"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            format_frame, text="僅 Markdown", 
            variable=self.output_format_var, value="markdown"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            format_frame, text="兩者都要", 
            variable=self.output_format_var, value="both"
        ).pack(side=tk.LEFT, padx=5)
        
        # Markdown 處理選項
        md_frame = tk.Frame(frame)
        md_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(md_frame, text="Markdown 處理:").pack(side=tk.LEFT, padx=10)
        
        self.md_process_var = tk.StringVar(value="basic")
        
        tk.Radiobutton(
            md_frame, text="基本處理", 
            variable=self.md_process_var, value="basic"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            md_frame, text="使用 MarkItDown (推薦)", 
            variable=self.md_process_var, value="markitdown"
        ).pack(side=tk.LEFT, padx=5)
        
        # API Key 設置
        api_frame = tk.Frame(frame)
        api_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(api_frame, text="OpenAI API Key:").pack(side=tk.LEFT, padx=10)
        self.api_key_entry = tk.Entry(api_frame, width=50)
        self.api_key_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        # 從環境變數獲取 API Key
        api_key_env = os.environ.get("OPENAI_API_KEY", "")
        if api_key_env:
            self.api_key_entry.insert(0, api_key_env)
            
        # 提示文字
        tip_label = tk.Label(
            frame, 
            text="注意: 使用 MarkItDown 可以更好地提取幻燈片中的文字，但需要安裝 markitdown 套件",
            font=("Arial", 10), fg="gray"
        )
        tip_label.pack(pady=5)
        
        # 標題設置
        title_frame = tk.Frame(frame)
        title_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(title_frame, text="文檔標題:").pack(side=tk.LEFT, padx=10)
        self.title_entry = tk.Entry(title_frame, width=50)
        self.title_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        self.title_entry.insert(0, "視頻捕獲的幻燈片")
        
        # 狀態顯示
        self.process_status_var = tk.StringVar(value="準備就緒")
        status_label = tk.Label(
            frame, textvariable=self.process_status_var,
            font=("Arial", 10), fg="blue"
        )
        status_label.pack(pady=10)
        
        # 進度條
        self.process_progress = ttk.Progressbar(
            frame, orient="horizontal", length=300, mode="indeterminate"
        )
        self.process_progress.pack(pady=10)
        
        # 處理按鈕
        self.process_btn = tk.Button(
            frame, text="處理幻燈片", command=self.process_slides,
            bg="#4CAF50", fg="white", height=2, width=20
        )
        self.process_btn.pack(pady=20)
    
    def browse_process_folder(self):
        """瀏覽並選擇幻燈片處理的源文件夾"""
        folder_path = filedialog.askdirectory(
            initialdir=".",
            title="選擇包含幻燈片的文件夾"
        )
        if folder_path:
            self.process_folder_entry.delete(0, tk.END)
            self.process_folder_entry.insert(0, folder_path)
    
    def process_slides(self):
        """處理幻燈片按鈕點擊處理"""
        folder = self.process_folder_entry.get()
        output_format = self.output_format_var.get()
        md_process = self.md_process_var.get()
        api_key = self.api_key_entry.get()
        title = self.title_entry.get()
        
        if not folder:
            messagebox.showwarning("警告", "請選擇幻燈片文件夾")
            return
        
        if not os.path.exists(folder):
            messagebox.showerror("錯誤", f"文件夾不存在: {folder}")
            return
            
        # 確保選擇的是目錄而不是文件
        if not os.path.isdir(folder):
            messagebox.showerror("錯誤", f"選擇的路徑不是目錄: {folder}")
            return
            
        # 檢查文件夾中是否有圖片
        has_images = False
        for filename in os.listdir(folder):
            if filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                has_images = True
                break
                
        if not has_images:
            messagebox.showwarning("警告", "選定的文件夾中未找到圖片")
            return
            
        # 開始處理（在背景線程中運行）
        self.process_status_var.set("正在處理幻燈片...")
        self.process_progress.start(10)
        self.process_btn.config(state=tk.DISABLED)
        
        def process_thread():
            # 處理結果
            results = []
            success = True
            
            # 根據選擇處理 PowerPoint
            if output_format in ["pptx", "both"]:
                ppt_success, ppt_result = generate_ppt_from_images(
                    folder, None, title
                )
                success = success and ppt_success
                if ppt_success:
                    results.append(f"PPT: {ppt_result}")
                else:
                    results.append(f"PPT 失敗: {ppt_result}")
            
            # 根據選擇處理 Markdown
            if output_format in ["markdown", "both"]:
                use_markitdown = (md_process == "markitdown")
                md_success, md_result = generate_markdown_from_images(
                    folder, None, title, use_markitdown, api_key
                )
                success = success and md_success
                if md_success:
                    results.append(f"Markdown: {md_result}")
                else:
                    results.append(f"Markdown 失敗: {md_result}")
            
            # 在主線程中更新 UI
            self.root.after(0, lambda: self.processing_completed(success, results))
        
        threading.Thread(target=process_thread).start()
    
    def processing_completed(self, success, results):
        """幻燈片處理完成後的處理"""
        self.process_progress.stop()
        self.process_btn.config(state=tk.NORMAL)
        
        if success:
            self.process_status_var.set("處理完成")
            
            # 顯示處理結果
            result_text = "\n".join(results)
            messagebox.showinfo("成功", f"幻燈片處理完成:\n{result_text}")
        else:
            self.process_status_var.set("處理失敗")
            
            # 顯示處理結果
            result_text = "\n".join(results)
            messagebox.showerror("錯誤", f"幻燈片處理失敗:\n{result_text}")

    def setup_web_capture_ui(self):
        """設置網頁捕獲界面"""
        frame = self.web_capture_frame
        
        # 頂部說明文字
        info_label = tk.Label(
            frame, 
            text="此功能可從截圖文件夾直接生成 PowerPoint 投影片",
            font=("Arial", 12)
        )
        info_label.pack(pady=20)
        
        # 截圖文件夾選擇區域
        folder_frame = tk.Frame(frame)
        folder_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(folder_frame, text="截圖文件夾:").pack(side=tk.LEFT, padx=10)
        self.screenshots_folder_entry = tk.Entry(folder_frame, width=50)
        self.screenshots_folder_entry.pack(
            side=tk.LEFT, padx=5, fill=tk.X, expand=True
        )
        
        self.screenshots_browse_btn = tk.Button(
            folder_frame, text="瀏覽...", 
            command=self.browse_screenshots_folder
        )
        self.screenshots_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # 輸出 PPT 文件選擇
        output_frame = tk.Frame(frame)
        output_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(output_frame, text="輸出 PPT 文件:").pack(side=tk.LEFT, padx=10)
        self.ppt_output_entry = tk.Entry(output_frame, width=50)
        self.ppt_output_entry.pack(
            side=tk.LEFT, padx=5, fill=tk.X, expand=True
        )
        
        self.ppt_browse_btn = tk.Button(
            output_frame, text="瀏覽...", 
            command=self.browse_ppt_output
        )
        self.ppt_browse_btn.pack(side=tk.LEFT, padx=5)
        
        # 標題設置
        title_frame = tk.Frame(frame)
        title_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(title_frame, text="PPT 標題:").pack(side=tk.LEFT, padx=10)
        self.ppt_title_entry = tk.Entry(title_frame, width=50)
        self.ppt_title_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        self.ppt_title_entry.insert(0, "截圖生成的投影片")
        
        # 投影片比例選擇
        ratio_frame = tk.Frame(frame)
        ratio_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(ratio_frame, text="投影片比例:").pack(side=tk.LEFT, padx=10)
        
        self.slide_ratio_var = tk.StringVar(value="16:9")
        
        tk.Radiobutton(
            ratio_frame, text="寬屏 (16:9)", 
            variable=self.slide_ratio_var, value="16:9"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            ratio_frame, text="標準 (4:3)", 
            variable=self.slide_ratio_var, value="4:3"
        ).pack(side=tk.LEFT, padx=5)
        
        # 圖片排序選項
        sort_frame = tk.Frame(frame)
        sort_frame.pack(fill=tk.X, pady=10)
        
        tk.Label(sort_frame, text="圖片排序:").pack(side=tk.LEFT, padx=10)
        
        self.sort_option_var = tk.StringVar(value="name")
        
        tk.Radiobutton(
            sort_frame, text="按文件名排序", 
            variable=self.sort_option_var, value="name"
        ).pack(side=tk.LEFT, padx=5)
        
        tk.Radiobutton(
            sort_frame, text="按修改時間排序", 
            variable=self.sort_option_var, value="time"
        ).pack(side=tk.LEFT, padx=5)
        
        # 圖片預覽框架
        preview_frame = tk.Frame(frame)
        preview_frame.pack(fill=tk.BOTH, expand=True, pady=10)
        
        # 預覽標籤
        tk.Label(
            preview_frame, text="圖片預覽 (顯示前 5 張):"
        ).pack(pady=5)
        
        # 預覽區域
        self.preview_canvas = tk.Canvas(
            preview_frame, bg="white", height=150
        )
        self.preview_canvas.pack(fill=tk.X, expand=True, padx=10)
        
        # 狀態顯示
        self.slide_make_status_var = tk.StringVar(value="準備就緒")
        status_label = tk.Label(
            frame, textvariable=self.slide_make_status_var,
            font=("Arial", 10), fg="blue"
        )
        status_label.pack(pady=10)
        
        # 生成按鈕
        self.make_ppt_btn = tk.Button(
            frame, text="生成 PowerPoint", 
            command=self.make_ppt_from_screenshots,
            bg="#4CAF50", fg="white", height=2, width=20
        )
        self.make_ppt_btn.pack(pady=20)
        
    def browse_screenshots_folder(self):
        """瀏覽並選擇截圖文件夾"""
        folder_path = filedialog.askdirectory(
            initialdir=".",
            title="選擇包含截圖的文件夾"
        )
        if folder_path:
            self.screenshots_folder_entry.delete(0, tk.END)
            self.screenshots_folder_entry.insert(0, folder_path)
            
            # 更新輸出 PPT 名稱
            folder_name = os.path.basename(folder_path)
            ppt_path = os.path.join(
                os.path.dirname(folder_path),
                f"{folder_name}.pptx"
            )
            
            self.ppt_output_entry.delete(0, tk.END)
            self.ppt_output_entry.insert(0, ppt_path)
            
            # 更新預覽
            self.update_preview()
            
    def browse_ppt_output(self):
        """選擇 PPT 輸出文件"""
        file_path = filedialog.asksaveasfilename(
            initialdir=".",
            title="保存 PowerPoint 文件",
            filetypes=[("PowerPoint 文件", "*.pptx")]
        )
        if file_path:
            if not file_path.lower().endswith(".pptx"):
                file_path += ".pptx"
                
            self.ppt_output_entry.delete(0, tk.END)
            self.ppt_output_entry.insert(0, file_path)
    
    def update_preview(self):
        """更新圖片預覽"""
        folder_path = self.screenshots_folder_entry.get()
        if not folder_path or not os.path.isdir(folder_path):
            return
            
        # 清除現有預覽
        self.preview_canvas.delete("all")
        
        # 獲取所有圖片文件
        image_files = self.get_image_files_sorted()
        
        if not image_files:
            self.preview_canvas.create_text(
                self.preview_canvas.winfo_width() // 2, 
                self.preview_canvas.winfo_height() // 2,
                text="找不到圖片文件",
                font=("Arial", 14)
            )
            return
            
        # 只顯示前 5 張
        preview_count = min(5, len(image_files))
        preview_images = []
        
        # 計算每張預覽圖的大小和位置
        canvas_width = self.preview_canvas.winfo_width()
        if canvas_width < 100:  # 如果畫布尚未渲染，使用預設寬度
            canvas_width = 500
            
        preview_width = (canvas_width - 10 * (preview_count + 1)) // preview_count
        preview_height = 130
        
        # 載入並顯示預覽圖
        for i in range(preview_count):
            try:
                # 打開圖片並調整大小
                img = Image.open(image_files[i])
                img.thumbnail((preview_width, preview_height), Image.LANCZOS)
                
                # 轉換為 PhotoImage
                photo = ImageTk.PhotoImage(img)
                preview_images.append(photo)  # 保持引用
                
                # 計算位置
                x = 10 + i * (preview_width + 10)
                
                # 在畫布上顯示圖片
                self.preview_canvas.create_image(x, 10, image=photo, anchor=tk.NW)
                
                # 顯示文件名
                filename = os.path.basename(image_files[i])
                if len(filename) > 15:
                    filename = filename[:12] + "..."
                
                self.preview_canvas.create_text(
                    x + preview_width // 2, 
                    preview_height + 20,
                    text=filename,
                    font=("Arial", 8)
                )
                
            except Exception as e:
                print(f"無法載入預覽圖: {e}")
        
        # 保持圖片引用，避免垃圾回收
        self.preview_canvas.preview_images = preview_images
    
    def get_image_files_sorted(self):
        """獲取排序後的圖片文件列表"""
        folder_path = self.screenshots_folder_entry.get()
        if not folder_path or not os.path.isdir(folder_path):
            return []
            
        # 獲取所有圖片文件
        image_files = []
        for filename in os.listdir(folder_path):
            if filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                image_files.append(os.path.join(folder_path, filename))
                
        # 根據選擇的排序方式進行排序
        sort_option = self.sort_option_var.get()
        
        if sort_option == "name":
            # 按文件名排序
            image_files.sort()
        elif sort_option == "time":
            # 按修改時間排序
            image_files.sort(key=os.path.getmtime)
            
        return image_files
    
    def make_ppt_from_screenshots(self):
        """從截圖生成 PowerPoint"""
        folder_path = self.screenshots_folder_entry.get()
        output_path = self.ppt_output_entry.get()
        title = self.ppt_title_entry.get()
        slide_ratio = self.slide_ratio_var.get()
        
        if not folder_path:
            messagebox.showwarning("警告", "請選擇截圖文件夾")
            return
            
        if not os.path.isdir(folder_path):
            messagebox.showerror("錯誤", f"文件夾不存在: {folder_path}")
            return
            
        # 獲取排序後的圖片文件
        image_files = self.get_image_files_sorted()
        
        if not image_files:
            messagebox.showwarning("警告", "選定的文件夾中未找到圖片")
            return
            
        # 確保輸出路徑有效
        if not output_path or output_path.strip() == "" or output_path.strip() == ".pptx":
            # 使用截圖文件夾名稱作為檔案名
            folder_name = os.path.basename(folder_path)
            if not folder_name or folder_name.strip() == "":
                folder_name = "投影片"  # 默認檔案名
                
            output_path = os.path.join(
                os.path.dirname(folder_path),
                f"{folder_name}.pptx"
            )
            self.ppt_output_entry.delete(0, tk.END)
            self.ppt_output_entry.insert(0, output_path)
        elif not output_path.lower().endswith(".pptx"):
            # 確保有正確的副檔名
            output_path += ".pptx"
            self.ppt_output_entry.delete(0, tk.END)
            self.ppt_output_entry.insert(0, output_path)
            
        # 開始生成（在背景線程中運行）
        self.slide_make_status_var.set("正在生成 PowerPoint...")
        self.make_ppt_btn.config(state=tk.DISABLED)
        
        def generate_thread():
            success, result = generate_ppt_from_images(
                folder_path, output_path, title, image_files, slide_ratio
            )
            
            # 在主線程中更新 UI
            self.root.after(
                0, lambda: self.ppt_generation_completed(success, result)
            )
        
        threading.Thread(target=generate_thread).start()
    
    def ppt_generation_completed(self, success, result):
        """PowerPoint 生成完成後的處理"""
        self.make_ppt_btn.config(state=tk.NORMAL)
        
        if success:
            self.slide_make_status_var.set("PowerPoint 生成成功")
            messagebox.showinfo("成功", f"PowerPoint 已成功生成: {result}")
        else:
            self.slide_make_status_var.set(f"PowerPoint 生成失敗: {result}")
            messagebox.showerror("錯誤", f"PowerPoint 生成失敗: {result}")

    def browse_slide_folder(self):
        """瀏覽並選擇幻燈片輸出文件夾"""
        folder_path = filedialog.askdirectory(
            initialdir=".",
            title="選擇保存幻燈片的文件夾"
        )
        if folder_path:
            self.slide_output_entry.delete(0, tk.END)
            self.slide_output_entry.insert(0, folder_path)

    # 新增更新函數 (確保綁定正確)
    def update_threshold_label(self, *args):
        value = self.threshold_var.get()
        self.threshold_label.config(text=f"{value:.2f}")
    
    def update_stability_label(self, *args):
        value = self.stability_var.get()
        self.stability_label.config(text=f"{value:.2f}")
    
    def update_interval_label(self, *args):
        value = self.interval_var.get()
        self.interval_label.config(text=f"{value:.2f}")
    
    def update_speed_label(self, *args):
        value = self.speed_var.get()
        self.speed_label.config(text=f"{value:.1f}x")

    # 新增預設設定應用函數
    def apply_preset(self, event):
        """套用預設設定"""
        preset = self.preset_var.get()
        
        if preset == "標準設定 (推薦)":
            self.threshold_var.set(0.85)  # 提高到 0.85
            self.stability_var.set(0.96)  # 提高到 0.96
            self.interval_var.set(1.0)    # 增加到 1.0 秒
            self.speed_var.set(1.5)       # 1.5倍速
            self.fast_mode_var.set(True)  # 啟用快速模式
            self.preset_tip.config(text="適合大多數影片，減少重複捕獲，中等加速")
            
        elif preset == "動畫較多的簡報":
            self.threshold_var.set(0.75)  # 提高到 0.75
            self.stability_var.set(0.98)  # 提高到 0.98
            self.interval_var.set(0.5)    # 保持 0.5 秒
            self.speed_var.set(1.0)       # 正常速度，動畫需要仔細捕獲
            self.fast_mode_var.set(False) # 關閉快速模式以更好地捕獲動畫
            self.preset_tip.config(text="適用於有過場動畫的簡報，正常速度")
            
        elif preset == "靜態畫面為主的簡報":
            self.threshold_var.set(0.90)  # 提高到 0.90
            self.stability_var.set(0.94)  # 稍微提高
            self.interval_var.set(2.0)    # 增加到 2.0 秒
            self.speed_var.set(2.0)       # 2倍速，靜態畫面可以快速處理
            self.fast_mode_var.set(True)  # 啟用快速模式
            self.preset_tip.config(text="適用於畫面變化少的簡報，2倍速加快處理")
            
        elif preset == "長影片優化":
            self.threshold_var.set(0.88)  # 提高到 0.88
            self.stability_var.set(0.96)  # 保持 0.96
            self.interval_var.set(2.5)    # 增加到 2.5 秒
            self.speed_var.set(3.0)       # 3倍速，長影片需要更快處理
            self.fast_mode_var.set(True)  # 啟用快速模式
            self.preset_tip.config(text="適用於長時間影片處理，3倍速大幅加快")
            
        elif preset == "捕獲較少的設定 (靈敏度低)":
            self.threshold_var.set(0.92)  # 更高的相似度閾值
            self.stability_var.set(0.98)  # 保持高穩定性要求
            self.interval_var.set(2.0)    # 更長的取樣間隔
            self.speed_var.set(2.5)       # 2.5倍速
            self.fast_mode_var.set(True)  # 啟用快速模式
            self.preset_tip.config(text="大幅減少捕獲數量，只保留明顯不同的投影片，2.5倍速")
            
        elif preset == "快速模式 (4倍速)":
            self.threshold_var.set(0.88)  # 適中的相似度閾值
            self.stability_var.set(0.95)  # 稍微降低穩定性要求
            self.interval_var.set(1.5)    # 1.5秒間隔
            self.speed_var.set(4.0)       # 4倍速，最快模式
            self.fast_mode_var.set(True)  # 啟用快速模式
            self.preset_tip.config(text="最快速模式，4倍速處理，適合快速預覽")

    # 新增停止捕獲功能
    def stop_capture(self):
        """停止捕獲過程"""
        if hasattr(self, 'capture_thread') and self.capture_thread.is_alive():
            self.stop_capture_flag = True
            self.stop_btn.config(state=tk.DISABLED)
            self.slide_status_var.set("停止請求已發送...")


def main():
    """主函數"""
    # 檢查依賴
    missing_packages = check_dependencies()
    
    if missing_packages:
        print(f"缺少以下依賴: {', '.join(missing_packages)}")
        
        # 特別提示 OpenCV
        if "opencv-python" in missing_packages:
            print("注意: OpenCV 是幻燈片捕獲功能的必要依賴")
            print("請確保安裝正確版本: pip install opencv-python")
        
        choice = input("是否自動安裝這些依賴？(y/n): ")
        
        if choice.lower() == 'y':
            if not install_dependencies(missing_packages):
                print("依賴安裝失敗！")
                print("\n🔧 解決方案:")
                print("1. 推薦使用 start.sh 腳本（自動處理虛擬環境）:")
                print("   chmod +x start.sh")
                print("   ./start.sh")
                print("\n2. 或手動創建虛擬環境:")
                print("   python3 -m venv venv")
                print("   source venv/bin/activate")
                print(f"   pip install {' '.join(missing_packages)}")
                
                # 特別提示 OpenCV 安裝問題
                if "opencv-python" in missing_packages:
                    print("\n💡 OpenCV 安裝提示:")
                    print("如果 opencv-python 安裝失敗，可嘗試:")
                    print("pip install opencv-python-headless")
                
                sys.exit(1)
        else:
            print("\n🔧 手動安裝步驟:")
            print("1. 推薦使用 start.sh 腳本:")
            print("   chmod +x start.sh")
            print("   ./start.sh")
            print("\n2. 或手動創建虛擬環境安裝:")
            print("   python3 -m venv venv")
            print("   source venv/bin/activate")
            print(f"   pip install {' '.join(missing_packages)}")
            
            # 特別提示 OpenCV
            if "opencv-python" in missing_packages:
                print("\n💡 OpenCV 安裝提示:")
                print("如果 opencv-python 安裝失敗，可嘗試:")
                print("pip install opencv-python-headless")
            
            sys.exit(1)
    
    try:
        # 創建並運行應用
        root = tk.Tk()
        app = VideoAudioProcessor(root)
        
        # 顯示使用提示
        messagebox.showinfo(
            "歡迎使用", 
            "歡迎使用視頻和音頻處理工具\n\n"
            "本工具提供以下功能：\n"
            "1. 從視頻文件中提取音頻\n"
            "2. 從視頻文件中捕獲幻燈片\n"
            "3. 將捕獲的幻燈片轉換為 PowerPoint 或 Markdown\n\n"
            "請選擇相應的標籤頁開始使用"
        )
        
        root.mainloop()
        
    except Exception as e:
        error_msg = traceback.format_exc()
        print(f"啟動失敗: {str(e)}\n{error_msg}")
        messagebox.showerror("啟動錯誤", f"程序啟動失敗:\n{str(e)}")


if __name__ == "__main__":
    main() 