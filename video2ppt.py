#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import cv2
import numpy as np
import os
import time
import argparse
import threading
from datetime import datetime
from PIL import Image, ImageTk
import tkinter as tk
from tkinter import filedialog, messagebox, Scale
from pptx import Presentation
from pptx.util import Inches
from skimage.metrics import structural_similarity as ssim
import pyautogui  # 添加pyautogui以支持螢幕擷取

class Video2PPT:
    def __init__(self, root):
        self.root = root
        self.root.title("Video2PPT - 影片投影片轉換器")
        self.root.geometry("1000x600")
        self.root.minsize(800, 500)
        
        # 預設值
        self.video_source = None
        self.cap = None
        self.is_running = False
        self.is_paused = False
        self.threshold = 0.95
        self.interval = 0.5
        self.output_file = "slides.pptx"
        self.slides_folder = "slides"
        self.slide_count = 0
        self.last_frame = None
        self.roi = None
        self.roi_selecting = False
        self.roi_start = None
        self.roi_end = None
        self.screen_capture_thread = None
        
        # 建立資料夾存放投影片
        if not os.path.exists(self.slides_folder):
            os.makedirs(self.slides_folder)
        
        self.setup_ui()
        
    def setup_ui(self):
        # 上方控制區域
        control_frame = tk.Frame(self.root)
        control_frame.pack(fill=tk.X, padx=10, pady=10)
        
        # 第一行控制項
        row1 = tk.Frame(control_frame)
        row1.pack(fill=tk.X, pady=5)
        
        tk.Label(row1, text="視頻來源:").pack(side=tk.LEFT, padx=5)
        self.source_var = tk.StringVar(value="文件")
        tk.Radiobutton(row1, text="文件", variable=self.source_var, value="文件", 
                      command=self.toggle_source).pack(side=tk.LEFT, padx=5)
        tk.Radiobutton(row1, text="URL", variable=self.source_var, value="URL", 
                      command=self.toggle_source).pack(side=tk.LEFT, padx=5)
        tk.Radiobutton(row1, text="螢幕", variable=self.source_var, value="螢幕", 
                      command=self.toggle_source).pack(side=tk.LEFT, padx=5)
        
        self.file_entry = tk.Entry(row1, width=50)
        self.file_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        self.browse_btn = tk.Button(row1, text="瀏覽", command=self.browse_file)
        self.browse_btn.pack(side=tk.LEFT, padx=5)
        
        # 第二行控制項
        row2 = tk.Frame(control_frame)
        row2.pack(fill=tk.X, pady=5)
        
        tk.Label(row2, text="相似度閾值:").pack(side=tk.LEFT, padx=5)
        self.threshold_scale = Scale(row2, from_=0.5, to=1.0, resolution=0.01, 
                                    orient=tk.HORIZONTAL, length=200)
        self.threshold_scale.set(self.threshold)
        self.threshold_scale.pack(side=tk.LEFT, padx=5)
        
        tk.Label(row2, text="檢測間隔(秒):").pack(side=tk.LEFT, padx=5)
        self.interval_scale = Scale(row2, from_=0.1, to=5.0, resolution=0.1, 
                                   orient=tk.HORIZONTAL, length=200)
        self.interval_scale.set(self.interval)
        self.interval_scale.pack(side=tk.LEFT, padx=5)
        
        tk.Label(row2, text="輸出文件:").pack(side=tk.LEFT, padx=5)
        self.output_entry = tk.Entry(row2, width=20)
        self.output_entry.insert(0, self.output_file)
        self.output_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        # 第三行 - 按鈕區
        row3 = tk.Frame(control_frame)
        row3.pack(fill=tk.X, pady=5)
        
        self.start_btn = tk.Button(row3, text="開始", command=self.start_capture, bg="#4CAF50", fg="white", width=10)
        self.start_btn.pack(side=tk.LEFT, padx=5)
        
        self.pause_btn = tk.Button(row3, text="暫停", command=self.pause_capture, state=tk.DISABLED, width=10)
        self.pause_btn.pack(side=tk.LEFT, padx=5)
        
        self.stop_btn = tk.Button(row3, text="停止", command=self.stop_capture, state=tk.DISABLED, width=10)
        self.stop_btn.pack(side=tk.LEFT, padx=5)
        
        self.reset_roi_btn = tk.Button(row3, text="重設區域", command=self.reset_roi, width=10)
        self.reset_roi_btn.pack(side=tk.LEFT, padx=5)
        
        self.generate_btn = tk.Button(row3, text="生成PPT", command=self.generate_ppt, width=10)
        self.generate_btn.pack(side=tk.LEFT, padx=5)
        
        # 右側狀態信息
        self.status_var = tk.StringVar(value="就緒")
        status_label = tk.Label(row3, textvariable=self.status_var, bd=1, relief=tk.SUNKEN, anchor=tk.W)
        status_label.pack(side=tk.RIGHT, padx=5, fill=tk.X, expand=True)
        
        # 視頻顯示區域
        self.canvas_frame = tk.Frame(self.root, bg="black")
        self.canvas_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        self.canvas = tk.Canvas(self.canvas_frame, bg="black")
        self.canvas.pack(fill=tk.BOTH, expand=True)
        
        # 綁定鼠標事件用於選擇ROI
        self.canvas.bind("<ButtonPress-1>", self.on_mouse_down)
        self.canvas.bind("<B1-Motion>", self.on_mouse_move)
        self.canvas.bind("<ButtonRelease-1>", self.on_mouse_up)
        
        # 下方日誌區域
        log_frame = tk.Frame(self.root)
        log_frame.pack(fill=tk.X, padx=10, pady=5)
        
        tk.Label(log_frame, text="執行日誌:").pack(anchor=tk.W)
        self.log_text = tk.Text(log_frame, height=5, width=50)
        self.log_text.pack(fill=tk.X, expand=True)
        
        scrollbar = tk.Scrollbar(self.log_text)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.log_text.config(yscrollcommand=scrollbar.set)
        scrollbar.config(command=self.log_text.yview)
        
        self.log("程序已啟動")
        self.log("請選擇視頻文件或輸入URL，並設置參數")
    
    def toggle_source(self):
        source = self.source_var.get()
        if source == "文件":
            self.browse_btn.config(state=tk.NORMAL)
            self.file_entry.config(state=tk.NORMAL)
            self.file_entry.delete(0, tk.END)
        elif source == "URL":
            self.browse_btn.config(state=tk.DISABLED)
            self.file_entry.config(state=tk.NORMAL)
            self.file_entry.delete(0, tk.END)
            self.file_entry.insert(0, "rtsp://")
        elif source == "螢幕":
            self.browse_btn.config(state=tk.DISABLED)
            self.file_entry.config(state=tk.DISABLED)
            self.file_entry.delete(0, tk.END)
            self.file_entry.insert(0, "正在使用螢幕擷取模式")
            # 顯示當前螢幕預覽
            self.show_screen_preview()
    
    def show_screen_preview(self):
        """顯示螢幕預覽"""
        try:
            # 獲取螢幕截圖
            screenshot = pyautogui.screenshot()
            
            # 轉換為numpy數組
            frame = np.array(screenshot)
            frame = cv2.cvtColor(frame, cv2.COLOR_RGB2BGR)
            
            # 顯示在畫布上
            self.display_frame(frame)
            
            self.log("已顯示螢幕預覽，請框選要監控的投影片區域")
            
            # 設置幀大小，用於ROI計算
            self.frame_size = (frame.shape[1], frame.shape[0])
        except Exception as e:
            self.log(f"無法獲取螢幕預覽: {str(e)}")
            messagebox.showerror("錯誤", f"無法獲取螢幕預覽: {str(e)}")
    
    def browse_file(self):
        file_path = filedialog.askopenfilename(
            filetypes=[("Video files", "*.mp4 *.avi *.mkv *.mov"), ("All files", "*.*")]
        )
        if file_path:
            self.file_entry.delete(0, tk.END)
            self.file_entry.insert(0, file_path)
    
    def log(self, message):
        timestamp = datetime.now().strftime("%H:%M:%S")
        self.log_text.insert(tk.END, f"[{timestamp}] {message}\n")
        self.log_text.see(tk.END)
    
    def on_mouse_down(self, event):
        if not self.is_running and not self.is_paused:
            self.roi_selecting = True
            self.roi_start = (event.x, event.y)
            self.canvas.delete("roi")
    
    def on_mouse_move(self, event):
        if self.roi_selecting:
            self.roi_end = (event.x, event.y)
            self.canvas.delete("roi_temp")
            self.canvas.create_rectangle(
                self.roi_start[0], self.roi_start[1],
                self.roi_end[0], self.roi_end[1],
                outline="red", width=2, tags="roi_temp"
            )
    
    def on_mouse_up(self, event):
        if self.roi_selecting:
            self.roi_selecting = False
            self.roi_end = (event.x, event.y)
            
            # 計算ROI
            x1, y1 = min(self.roi_start[0], self.roi_end[0]), min(self.roi_start[1], self.roi_end[1])
            x2, y2 = max(self.roi_start[0], self.roi_end[0]), max(self.roi_start[1], self.roi_end[1])
            
            # 保存實際像素範圍
            if hasattr(self, 'frame_size'):
                # 計算實際像素與顯示像素的比例
                ratio_x = self.frame_size[0] / self.display_size[0]
                ratio_y = self.frame_size[1] / self.display_size[1]
                
                # 計算實際像素位置
                self.roi = (
                    int(x1 * ratio_x), int(y1 * ratio_y),
                    int(x2 * ratio_x), int(y2 * ratio_y)
                )
            else:
                self.roi = (x1, y1, x2, y2)
            
            self.canvas.delete("roi_temp")
            self.canvas.create_rectangle(
                x1, y1, x2, y2,
                outline="green", width=2, tags="roi"
            )
            
            self.log(f"已選擇區域: {self.roi}")
    
    def reset_roi(self):
        self.roi = None
        self.canvas.delete("roi")
        self.log("已重設選擇區域")
    
    def start_capture(self):
        # 取得參數
        self.threshold = float(self.threshold_scale.get())
        self.interval = float(self.interval_scale.get())
        self.output_file = self.output_entry.get()
        
        # 檢查輸入
        source = self.source_var.get()
        
        if source == "文件" or source == "URL":
            path_or_url = self.file_entry.get().strip()
            
            if not path_or_url:
                messagebox.showerror("錯誤", "請選擇視頻文件或輸入URL")
                return
            
            if source == "文件" and not os.path.exists(path_or_url):
                messagebox.showerror("錯誤", "所選擇的文件不存在")
                return
                
            # 打開視頻
            self.video_source = path_or_url
            self.cap = cv2.VideoCapture(self.video_source)
            
            if not self.cap.isOpened():
                messagebox.showerror("錯誤", "無法打開視頻源")
                return
                
            # 獲取視頻尺寸
            ret, frame = self.cap.read()
            if not ret:
                messagebox.showerror("錯誤", "無法讀取視頻幀")
                return
                
            self.frame_size = (frame.shape[1], frame.shape[0])
        
        elif source == "螢幕":
            # 確保已選擇ROI
            if not self.roi:
                messagebox.showerror("錯誤", "請先框選要監控的投影片區域")
                return
            
            # 設置螢幕捕獲
            self.video_source = "screen"
            # 不需要self.cap
            
        # 更新UI狀態
        self.start_btn.config(state=tk.DISABLED)
        self.pause_btn.config(state=tk.NORMAL)
        self.stop_btn.config(state=tk.NORMAL)
        self.is_running = True
        self.is_paused = False
        
        self.log(f"開始捕獲: {self.video_source}")
        self.log(f"相似度閾值: {self.threshold}, 間隔: {self.interval}秒")
        
        # 開始處理
        if source == "螢幕":
            # 啟動螢幕捕獲線程
            self.screen_capture_thread = threading.Thread(target=self.capture_screen)
            self.screen_capture_thread.daemon = True
            self.screen_capture_thread.start()
        else:
            # 開始處理視頻
            self.process_video()
    
    def capture_screen(self):
        """螢幕捕獲線程"""
        try:
            self.log("開始螢幕擷取")
            x1, y1, x2, y2 = self.roi
            width = x2 - x1
            height = y2 - y1
            
            last_screenshot = None
            
            while self.is_running and not self.is_paused:
                # 獲取螢幕截圖
                screenshot = pyautogui.screenshot(region=(x1, y1, width, height))
                
                # 轉換為OpenCV格式
                frame = np.array(screenshot)
                frame = cv2.cvtColor(frame, cv2.COLOR_RGB2BGR)
                
                # 在主線程中顯示截圖
                self.root.after(0, lambda: self.display_capture(frame))
                
                # 檢查是否有變化
                if last_screenshot is not None:
                    # 轉換為灰度圖像以進行比較
                    gray1 = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
                    gray2 = cv2.cvtColor(last_screenshot, cv2.COLOR_BGR2GRAY)
                    
                    # 計算相似度
                    similarity = ssim(gray1, gray2)
                    
                    # 更新狀態信息
                    self.root.after(0, lambda s=similarity: self.status_var.set(f"相似度: {s:.4f}"))
                    
                    # 如果幀間差異足夠大，保存截圖
                    if similarity < self.threshold:
                        # 在主線程中保存投影片
                        self.root.after(0, lambda f=frame: self.save_slide(f))
                        last_screenshot = frame
                else:
                    last_screenshot = frame
                    # 保存第一幀
                    self.root.after(0, lambda f=frame: self.save_slide(f))
                
                # 睡眠一段時間
                time.sleep(self.interval)
                
        except Exception as e:
            self.log(f"螢幕擷取出錯: {str(e)}")
            self.root.after(0, lambda: self.stop_capture())
    
    def display_capture(self, frame):
        """在主線程中顯示捕獲的幀"""
        self.display_frame(frame)
    
    def process_video(self):
        if not self.is_running or self.is_paused:
            return
            
        ret, frame = self.cap.read()
        if not ret:
            self.log("視頻結束")
            self.stop_capture()
            return
        
        # 顯示當前幀
        self.display_frame(frame)
        
        # 如果定義了ROI，處理這部分
        if self.roi:
            x1, y1, x2, y2 = self.roi
            roi_frame = frame[y1:y2, x1:x2]
            
            # 檢查是否有變化
            if self.last_frame is not None:
                # 確保兩個幀大小相同
                if roi_frame.shape == self.last_frame.shape:
                    # 轉換為灰度圖像以進行比較
                    gray1 = cv2.cvtColor(roi_frame, cv2.COLOR_BGR2GRAY)
                    gray2 = cv2.cvtColor(self.last_frame, cv2.COLOR_BGR2GRAY)
                    
                    # 計算相似度
                    similarity = ssim(gray1, gray2)
                    
                    status_text = f"相似度: {similarity:.4f}"
                    self.status_var.set(status_text)
                    
                    # 如果幀間差異足夠大，保存截圖
                    if similarity < self.threshold:
                        self.save_slide(roi_frame)
                        self.last_frame = roi_frame
                else:
                    self.last_frame = roi_frame
            else:
                self.last_frame = roi_frame
                self.save_slide(roi_frame)
        
        # 設置下一次處理
        self.root.after(int(self.interval * 1000), self.process_video)
        
    def display_frame(self, frame):
        # 調整大小以適應畫布
        canvas_width = self.canvas.winfo_width()
        canvas_height = self.canvas.winfo_height()
        
        if canvas_width <= 1 or canvas_height <= 1:
            canvas_width = self.canvas_frame.winfo_width()
            canvas_height = self.canvas_frame.winfo_height()
        
        # 保持寬高比
        scale = min(canvas_width/frame.shape[1], canvas_height/frame.shape[0])
        new_width = int(frame.shape[1] * scale)
        new_height = int(frame.shape[0] * scale)
        
        self.display_size = (new_width, new_height)
        
        # 調整大小並轉換為PIL圖像
        frame_resized = cv2.resize(frame, (new_width, new_height))
        frame_rgb = cv2.cvtColor(frame_resized, cv2.COLOR_BGR2RGB)
        pil_img = Image.fromarray(frame_rgb)
        
        # 轉換為Tkinter可用的格式
        img_tk = ImageTk.PhotoImage(image=pil_img)
        
        # 保持引用以防垃圾回收
        self.current_image = img_tk
        
        # 更新畫布
        self.canvas.delete("frame")
        self.canvas.create_image(
            canvas_width//2, canvas_height//2,
            image=img_tk, anchor=tk.CENTER, tags="frame"
        )
        
        # 如果有ROI，重新繪製
        if self.roi and not self.roi_selecting:
            x1, y1, x2, y2 = self.roi
            # 將實際像素位置轉換為顯示像素位置
            ratio_x = new_width / self.frame_size[0]
            ratio_y = new_height / self.frame_size[1]
            display_x1 = int(x1 * ratio_x)
            display_y1 = int(y1 * ratio_y)
            display_x2 = int(x2 * ratio_x)
            display_y2 = int(y2 * ratio_y)
            
            # 計算偏移量以居中顯示
            offset_x = (canvas_width - new_width) // 2
            offset_y = (canvas_height - new_height) // 2
            
            self.canvas.delete("roi")
            self.canvas.create_rectangle(
                display_x1 + offset_x, display_y1 + offset_y,
                display_x2 + offset_x, display_y2 + offset_y,
                outline="green", width=2, tags="roi"
            )
    
    def save_slide(self, frame):
        # 保存投影片截圖
        slide_filename = os.path.join(self.slides_folder, f"slide_{self.slide_count:03d}.png")
        cv2.imwrite(slide_filename, frame)
        
        self.slide_count += 1
        self.log(f"保存投影片 {self.slide_count}: {slide_filename}")
    
    def pause_capture(self):
        if self.is_running:
            if self.is_paused:
                self.is_paused = False
                self.pause_btn.config(text="暫停")
                self.log("繼續捕獲")
                
                # 如果是視頻模式，繼續處理
                if self.source_var.get() != "螢幕":
                    self.process_video()
                # 如果是螢幕模式，重新啟動螢幕捕獲線程
                else:
                    self.screen_capture_thread = threading.Thread(target=self.capture_screen)
                    self.screen_capture_thread.daemon = True
                    self.screen_capture_thread.start()
            else:
                self.is_paused = True
                self.pause_btn.config(text="繼續")
                self.log("暫停捕獲")
    
    def stop_capture(self):
        self.is_running = False
        self.is_paused = False
        
        if self.cap:
            self.cap.release()
            self.cap = None
        
        self.start_btn.config(state=tk.NORMAL)
        self.pause_btn.config(state=tk.DISABLED)
        self.stop_btn.config(state=tk.DISABLED)
        
        self.log("停止捕獲")
        
        if self.slide_count > 0:
            self.log(f"總共捕獲 {self.slide_count} 張投影片")
            messagebox.showinfo("完成", f"已捕獲 {self.slide_count} 張投影片。\n您可以點擊「生成PPT」按鈕創建簡報。")
    
    def generate_ppt(self):
        if self.slide_count == 0:
            messagebox.showwarning("警告", "沒有捕獲的投影片。請先運行視頻捕獲。")
            return
            
        try:
            # 創建PPT
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]  # 選擇空白佈局
            
            # 為每張圖片創建一張幻燈片
            for i in range(self.slide_count):
                slide_filename = os.path.join(self.slides_folder, f"slide_{i:03d}.png")
                
                if os.path.exists(slide_filename):
                    slide = prs.slides.add_slide(slide_layout)
                    # 添加圖片並調整大小以適應幻燈片
                    img_path = slide_filename
                    slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            
            # 保存PPT
            prs.save(self.output_file)
            self.log(f"已成功生成PPT文件: {self.output_file}")
            messagebox.showinfo("成功", f"已成功生成PPT文件: {self.output_file}")
            
        except Exception as e:
            self.log(f"生成PPT時出錯: {str(e)}")
            messagebox.showerror("錯誤", f"生成PPT時出錯: {str(e)}")
    
    def on_closing(self):
        if self.is_running:
            self.stop_capture()
        self.root.destroy()

def main():
    parser = argparse.ArgumentParser(description="Video2PPT - 視頻投影片自動捕獲器")
    parser.add_argument("--threshold", type=float, default=0.95, help="相似度閾值 (0.5-1.0)")
    parser.add_argument("--interval", type=float, default=0.5, help="檢測間隔(秒)")
    parser.add_argument("--output", type=str, default="slides.pptx", help="輸出PPT文件名")
    args = parser.parse_args()
    
    root = tk.Tk()
    app = Video2PPT(root)
    
    # 應用命令行參數
    app.threshold = args.threshold
    app.threshold_scale.set(args.threshold)
    app.interval = args.interval
    app.interval_scale.set(args.interval)
    app.output_file = args.output
    app.output_entry.delete(0, tk.END)
    app.output_entry.insert(0, args.output)
    
    root.protocol("WM_DELETE_WINDOW", app.on_closing)
    root.mainloop()

if __name__ == "__main__":
    main() 