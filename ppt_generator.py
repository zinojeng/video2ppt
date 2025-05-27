#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
PowerPoint生成工具：將截圖轉換為PowerPoint文件
"""

import os
import argparse
from pptx import Presentation
from pptx.util import Inches

def create_ppt_from_images(image_folder, output_file="presentation.pptx", title=None):
    """
    將一個資料夾中的圖片轉換為PowerPoint演示文稿
    
    參數:
        image_folder (str): 包含圖片的資料夾路徑
        output_file (str): 輸出的PowerPoint文件路徑
        title (str): 演示文稿標題（可選）
    
    返回:
        bool: 成功返回True，失敗返回False
    """
    try:
        # 確保輸出目錄存在
        output_dir = os.path.dirname(output_file)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        # 創建PowerPoint演示文稿
        prs = Presentation()
        
        # 如果指定了標題，添加標題幻燈片
        if title:
            title_slide_layout = prs.slide_layouts[0]  # 標題佈局
            slide = prs.slides.add_slide(title_slide_layout)
            title_placeholder = slide.shapes.title
            subtitle_placeholder = slide.placeholders[1]
            title_placeholder.text = title
            subtitle_placeholder.text = f"自動生成 - {len(os.listdir(image_folder))} 張投影片"
        
        # 使用空白佈局添加圖片
        blank_slide_layout = prs.slide_layouts[6]  # 空白佈局
        
        # 獲取圖片文件列表並排序
        image_files = [f for f in os.listdir(image_folder) 
                      if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
        image_files.sort()  # 按文件名排序
        
        # 添加每張圖片到幻燈片
        for img_file in image_files:
            img_path = os.path.join(image_folder, img_file)
            
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加圖片並調整大小以適合幻燈片
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), 
                                    width=prs.slide_width, 
                                    height=prs.slide_height)
            
            # 可選：添加頁碼
            # slide_number = len(prs.slides)
            # slide.shapes.add_textbox(Inches(9), Inches(6.5), Inches(1), Inches(0.5)).text = f"{slide_number}"
        
        # 保存PowerPoint文件
        prs.save(output_file)
        print(f"PowerPoint已成功創建並保存為: {output_file}")
        print(f"總共添加了 {len(image_files)} 張投影片")
        
        return True
        
    except Exception as e:
        print(f"創建PowerPoint時出錯: {str(e)}")
        return False

def main():
    # 解析命令行參數
    parser = argparse.ArgumentParser(description='將圖片轉換為PowerPoint演示文稿')
    parser.add_argument('image_folder', help='包含圖片的資料夾路徑')
    parser.add_argument('--output', '-o', default='presentation.pptx', help='輸出的PowerPoint文件路徑')
    parser.add_argument('--title', '-t', help='演示文稿標題（可選）')
    
    args = parser.parse_args()
    
    # 創建PowerPoint
    create_ppt_from_images(args.image_folder, args.output, args.title)

if __name__ == "__main__":
    main() 