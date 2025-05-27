#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
MarkItDown PPT 轉換工具

此工具用於將圖片檔案轉換為 Markdown 文件，然後生成 PPT 簡報檔案。
使用 MarkItDown 功能處理圖片內容，並可選用 AI 分析增強圖片描述。
"""

import os
import sys
import logging
import argparse
import glob
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
import time

# 設定日誌
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger("markitdown-ppt")


def convert_images_to_markdown(
    image_paths: List[str],
    output_file: str,
    title: str = "圖片集合",
    use_llm: bool = True,
    api_key: Optional[str] = None,
    model: str = "gpt-4o-mini",
    provider: str = "openai"
) -> Tuple[bool, str, Dict[str, Any]]:
    """
    將多個圖片檔案轉換為單一 Markdown 檔案
    
    Args:
        image_paths (List[str]): 圖片檔案路徑列表
        output_file (str): 輸出 Markdown 檔案路徑
        title (str): Markdown 文件標題
        use_llm (bool): 是否使用 LLM 處理圖片
        api_key (Optional[str]): API Key (OpenAI 或 Google)
        model (str): 使用的模型名稱
        provider (str): API 提供者，可為 'openai' 或 'gemini'
        
    Returns:
        Tuple[bool, str, Dict]: (是否成功, 輸出檔案路徑, 處理資訊)
    """
    try:
        # 嘗試導入必要模組
        try:
            from markitdown import MarkItDown
            
            # 根據提供者導入不同的 API 套件
            if provider.lower() == "gemini":
                try:
                    import google.generativeai as genai
                except ImportError:
                    logger.error("找不到 google.generativeai 模組，請確保已安裝 google-genai 套件")
                    return False, "", {"error": "缺少 Google Generative AI 模組"}
            else:  # 默認為 OpenAI
                try:
                    from openai import OpenAI, AuthenticationError
                except ImportError:
                    logger.error("找不到 openai 模組，請確保已安裝相關套件")
                    return False, "", {"error": "缺少 OpenAI 模組"}
            
        except ImportError:
            logger.error("找不到 markitdown 模組，請確保已安裝相關套件")
            return False, "", {"error": "缺少必要模組"}
            
        logger.info(f"將 {len(image_paths)} 張圖片轉換為 Markdown")
        
        # 檢查圖片列表
        if not image_paths:
            logger.error("圖片列表為空")
            return False, "", {"error": "圖片列表為空"}
            
        # 過濾有效的圖片檔案
        valid_images = []
        for img_path in image_paths:
            if os.path.exists(img_path):
                valid_images.append(img_path)
            else:
                logger.warning(f"找不到圖片: {img_path}")
                
        if not valid_images:
            logger.error("沒有有效的圖片檔案")
            return False, "", {"error": "沒有有效的圖片檔案"}
        
        # 生成初始 Markdown 文本
        md_content = f"# {title}\n\n"
        
        # 建立 MarkItDown 實例 (不使用內建 LLM，避免兼容性問題)
        md_kwargs = {"enable_plugins": True}
        llm_info = {}
        current_api_key = None
        
        if use_llm:
            logger.info(f"嘗試啟用 LLM ({model}) 進行處理，提供者: {provider}...")
            
            # 根據提供者獲取適合的 API 金鑰
            if provider.lower() == "gemini":
                current_api_key = api_key or os.environ.get("GOOGLE_API_KEY")
                
                if current_api_key:
                    try:
                        # 初始化 Google Generative AI
                        genai.configure(api_key=current_api_key)
                        
                        # 測試 API Key 是否有效
                        try:
                            # 簡單測試，列出可用模型
                            models = list(genai.list_models())
                            logger.info("Google API Key 驗證成功。")
                            llm_info["status"] = "啟用成功"
                            llm_info["model"] = model
                            llm_info["provider"] = "gemini"
                        except Exception as e:
                            logger.error(f"Google API Key 驗證失敗: {e}")
                            llm_info["status"] = "API Key 驗證失敗"
                            use_llm = False
                    except Exception as e:
                        logger.error(f"初始化 Gemini client 時發生錯誤: {e}")
                        llm_info["status"] = f"初始化錯誤: {str(e)}"
                        use_llm = False
                else:
                    logger.warning("未提供 Gemini API Key，無法使用 LLM 處理圖片。")
                    llm_info["status"] = "未提供 Gemini API Key"
                    use_llm = False
            else:  # 默認為 OpenAI
                current_api_key = api_key or os.environ.get("OPENAI_API_KEY")
                
                if not current_api_key:
                    logger.warning("未提供 OpenAI API Key，無法使用 LLM 處理圖片。")
                    llm_info["status"] = "未提供 OpenAI API Key"
                    use_llm = False
                else:
                    try:
                        # 初始化 OpenAI 客戶端
                        from openai import OpenAI, AuthenticationError
                        llm_client = OpenAI(api_key=current_api_key)
                        # 執行一個簡單的測試呼叫來驗證金鑰
                        llm_client.models.list() 
                        logger.info("OpenAI API Key 驗證成功。")
                        llm_info["status"] = "啟用成功"
                        llm_info["model"] = model
                        llm_info["provider"] = "openai"
                    except Exception as e:
                        logger.error(f"初始化 OpenAI client 時發生錯誤: {e}")
                        llm_info["status"] = f"初始化錯誤: {str(e)}"
                        use_llm = False
                    
        # 使用基本的 MarkItDown 實例，不傳遞 LLM 參數避免兼容性問題
        md = MarkItDown(**md_kwargs)
        
        # 處理每個圖片
        successful_conversions = 0
        for img_path in valid_images:
            try:
                img_relpath = os.path.basename(img_path)
                logger.info(f"處理圖片: {img_relpath}")
                
                # 轉換圖片
                result = md.convert(img_path)
                
                if result and result.text_content:
                    md_content += f"## 圖片：{img_relpath}\n\n"
                    md_content += result.text_content + "\n\n"
                    successful_conversions += 1
                else:
                    logger.warning(f"無法轉換圖片: {img_relpath}")
                    # 添加簡單的圖片標記
                    md_content += f"## 圖片：{img_relpath}\n\n"
                    md_content += f"![{img_relpath}]({img_path})\n\n"
            except Exception as e:
                logger.warning(f"處理圖片 {img_path} 時出錯: {e}")
                # 添加簡單的圖片標記
                md_content += f"## 圖片：{img_relpath}\n\n"
                md_content += f"![{img_relpath}]({img_path})\n\n"
                
        # 寫入輸出檔案
        try:
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(md_content)
            logger.info(f"已將 {successful_conversions} 張圖片轉換並寫入 {output_file}")
            
            # 嘗試使用 image_analyzer 增強圖片描述
            if use_llm and os.path.exists(output_file):
                try:
                    # 嘗試導入 image_analyzer
                    try:
                        from image_analyzer import enhance_markdown_with_image_analysis
                        
                        # 讀取剛生成的 Markdown 檔案
                        with open(output_file, 'r', encoding='utf-8') as f:
                            original_md = f.read()
                        
                        # 增強 Markdown 內容
                        enhanced_md, stats = enhance_markdown_with_image_analysis(
                            markdown_text=original_md,
                            base_dir=os.path.dirname(output_file),
                            api_key=current_api_key,
                            model=model,
                            provider=provider
                        )
                        
                        # 寫入增強後的內容
                        with open(output_file, 'w', encoding='utf-8') as f:
                            f.write(enhanced_md)
                        
                        logger.info(f"已增強 {stats['images_processed']} 張圖片的描述")
                    except ImportError:
                        logger.warning("找不到 image_analyzer 模組，無法增強圖片描述")
                except Exception as e:
                    logger.warning(f"增強圖片描述時出錯: {e}")
            
            result_info = {
                "success": True,
                "total_images": len(valid_images),
                "converted_images": successful_conversions,
                "output_file": output_file,
                "llm": llm_info
            }
            
            return True, output_file, result_info
            
        except Exception as e:
            logger.error(f"寫入輸出檔案時發生錯誤: {e}")
            return False, "", {"error": f"寫入輸出檔案時發生錯誤: {e}"}
        
    except Exception as e:
        logger.error(f"轉換圖片集合時發生錯誤: {e}")
        import traceback
        error_details = traceback.format_exc()
        logger.error(error_details)
        return False, "", {"error": str(e), "details": error_details}


def generate_ppt_from_markdown(
    markdown_file: str,
    output_file: str,
    template_file: Optional[str] = None
) -> bool:
    """
    從 Markdown 檔案生成 PowerPoint 簡報
    
    Args:
        markdown_file (str): Markdown 檔案路徑
        output_file (str): 輸出 PPT 檔案路徑
        template_file (Optional[str]): PPT 模板檔案路徑
        
    Returns:
        bool: 是否成功生成 PPT
    """
    try:
        # 檢查必要檔案
        if not os.path.exists(markdown_file):
            logger.error(f"找不到 Markdown 檔案: {markdown_file}")
            return False
            
        logger.info(f"從 Markdown 檔案生成 PPT: {markdown_file}")
        
        # 嘗試導入 python-pptx
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            import markdown
            from bs4 import BeautifulSoup
        except ImportError:
            logger.error("找不到必要的套件，請確保已安裝 python-pptx, markdown, beautifulsoup4")
            return False
            
        # 讀取 Markdown 內容
        with open(markdown_file, 'r', encoding='utf-8') as f:
            md_content = f.read()
            
        # 將 Markdown 轉換為 HTML
        html_content = markdown.markdown(md_content, extensions=['tables', 'fenced_code'])
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # 創建或載入 PPT
        if template_file and os.path.exists(template_file):
            prs = Presentation(template_file)
            logger.info(f"使用模板: {template_file}")
        else:
            prs = Presentation()
            logger.info("使用預設模板")
            
        # 找出所有標題和段落
        slides_content = []
        current_title = "簡報"
        current_content = []
        current_image = None
        
        # 處理每個 HTML 元素
        for element in soup.find_all(['h1', 'h2', 'p', 'img']):
            # 處理標題
            if element.name == 'h1':
                # 如果有內容，儲存前一個幻燈片
                if current_content or current_image:
                    slides_content.append((current_title, current_content, current_image))
                    current_content = []
                    current_image = None
                    
                # 新幻燈片標題
                current_title = element.text.strip()
                
            elif element.name == 'h2':
                # 如果有內容，儲存前一個幻燈片
                if current_content or current_image:
                    slides_content.append((current_title, current_content, current_image))
                    current_content = []
                    current_image = None
                    
                # 新幻燈片子標題
                current_title = element.text.strip()
                
            # 處理圖片
            elif element.name == 'img':
                img_src = element.get('src')
                if img_src:
                    if os.path.exists(img_src):
                        current_image = img_src
                    elif os.path.exists(os.path.join(os.path.dirname(markdown_file), img_src)):
                        current_image = os.path.join(os.path.dirname(markdown_file), img_src)
                    else:
                        logger.warning(f"找不到圖片: {img_src}")
                
            # 處理段落
            elif element.name == 'p':
                text = element.text.strip()
                if text:
                    current_content.append(text)
        
        # 儲存最後一個幻燈片
        if current_content or current_image:
            slides_content.append((current_title, current_content, current_image))
        
        # 創建幻燈片
        for title, content, image in slides_content:
            # 選擇幻燈片版型
            if image:
                slide_layout = prs.slide_layouts[8]  # 標題和內容與圖片
            else:
                slide_layout = prs.slide_layouts[1]  # 標題和內容
                
            # 創建幻燈片
            slide = prs.slides.add_slide(slide_layout)
            
            # 添加標題
            if slide.shapes.title:
                slide.shapes.title.text = title
            
            # 添加內容
            if content:
                content_text = "\n".join(content)
                if slide.placeholders[1]:
                    slide.placeholders[1].text = content_text
            
            # 添加圖片
            if image and os.path.exists(image):
                # 找出圖片放置的位置
                if len(slide.placeholders) > 2:
                    # 使用預留位置
                    placeholder = slide.placeholders[2]
                    placeholder.insert_picture(image)
                else:
                    # 手動添加圖片
                    left = Inches(1)
                    top = Inches(2.5)
                    width = Inches(4)
                    slide.shapes.add_picture(image, left, top, width=width)
        
        # 保存 PPT
        prs.save(output_file)
        logger.info(f"PPT 已生成: {output_file}")
        
        return True
        
    except Exception as e:
        logger.error(f"生成 PPT 時發生錯誤: {e}")
        import traceback
        error_details = traceback.format_exc()
        logger.error(error_details)
        return False


def process_images_to_ppt(
    image_dir: str,
    output_ppt: str,
    markdown_file: Optional[str] = None,
    title: str = "圖片簡報",
    use_llm: bool = True,
    api_key: Optional[str] = None,
    model: str = "gpt-4o-mini",
    provider: str = "openai",
    template_file: Optional[str] = None
) -> bool:
    """
    處理資料夾中的圖片，轉換為 Markdown 並生成 PPT
    
    Args:
        image_dir (str): 圖片資料夾路徑
        output_ppt (str): 輸出 PPT 檔案路徑
        markdown_file (Optional[str]): 輸出 Markdown 檔案路徑，如果為 None 則自動生成
        title (str): 簡報標題
        use_llm (bool): 是否使用 LLM 處理圖片
        api_key (Optional[str]): OpenAI 或 Google API Key
        model (str): 使用的模型名稱
        provider (str): API 提供者，可為 'openai' 或 'gemini'
        template_file (Optional[str]): PPT 模板檔案路徑
        
    Returns:
        bool: 是否成功生成 PPT
    """
    try:
        # 檢查圖片資料夾
        if not os.path.isdir(image_dir):
            logger.error(f"找不到圖片資料夾: {image_dir}")
            return False
            
        # 獲取圖片檔案列表
        image_files = []
        for ext in ['*.jpg', '*.jpeg', '*.png', '*.gif', '*.bmp']:
            image_files.extend(glob.glob(os.path.join(image_dir, ext)))
            image_files.extend(glob.glob(os.path.join(image_dir, ext.upper())))
            
        if not image_files:
            logger.error(f"資料夾中沒有圖片檔案: {image_dir}")
            return False
            
        # 按檔名排序
        image_files.sort()
        
        # 確定輸出 Markdown 檔案路徑
        if not markdown_file:
            markdown_file = os.path.join(
                os.path.dirname(output_ppt),
                f"{os.path.splitext(os.path.basename(output_ppt))[0]}.md"
            )
            
        # 轉換圖片為 Markdown
        logger.info(f"處理 {len(image_files)} 張圖片")
        success, md_file, info = convert_images_to_markdown(
            image_paths=image_files,
            output_file=markdown_file,
            title=title,
            use_llm=use_llm,
            api_key=api_key,
            model=model,
            provider=provider
        )
        
        if not success:
            logger.error(f"轉換圖片為 Markdown 失敗: {info.get('error', '未知錯誤')}")
            return False
            
        # 從 Markdown 生成 PPT
        logger.info(f"從 Markdown 生成 PPT: {md_file}")
        ppt_success = generate_ppt_from_markdown(
            markdown_file=md_file,
            output_file=output_ppt,
            template_file=template_file
        )
        
        if not ppt_success:
            logger.error("生成 PPT 失敗")
            return False
            
        logger.info(f"成功生成 PPT: {output_ppt}")
        return True
        
    except Exception as e:
        logger.error(f"處理圖片生成 PPT 時發生錯誤: {e}")
        import traceback
        error_details = traceback.format_exc()
        logger.error(error_details)
        return False


def process_captured_slides(
    slides_folder: str, 
    output_format: str = "markdown", 
    api_key: Optional[str] = None, 
    model: str = "gpt-4o-mini",
    provider: str = "openai"
) -> bool:
    """
    處理已捕獲的投影片，生成 Markdown 或 PowerPoint 或兩者都生成
    
    Args:
        slides_folder (str): 包含投影片圖片的資料夾路徑
        output_format (str): 輸出格式 ("markdown", "pptx", "both")
        api_key (Optional[str]): API Key (OpenAI 或 Google)
        model (str): 使用的模型名稱
        provider (str): API 提供者，可為 'openai' 或 'gemini'
        
    Returns:
        bool: 處理是否成功
    """
    logger.info(f"開始處理投影片: {slides_folder}, 格式: {output_format}")
    logger.info(f"使用 API 提供者: {provider}, 模型: {model}")
    
    # 檢查資料夾是否存在
    if not os.path.exists(slides_folder):
        logger.error(f"資料夾不存在: {slides_folder}")
        return False
    
    # 獲取所有圖片檔案
    image_files = []
    for filename in sorted(os.listdir(slides_folder)):
        if filename.lower().endswith(('.png', '.jpg', '.jpeg')):
            image_files.append(os.path.join(slides_folder, filename))
    
    if not image_files:
        logger.error(f"資料夾中沒有圖片檔案: {slides_folder}")
        return False
    
    success = True
    
    # 處理 Markdown
    if output_format in ["markdown", "both"]:
        output_md = os.path.join(
            os.path.dirname(slides_folder), 
            "slides_analysis.md"
        )
        
        md_success, _, _ = convert_images_to_markdown(
            image_paths=image_files,
            output_file=output_md,
            title="投影片內容分析",
            use_llm=(api_key is not None),
            api_key=api_key,
            model=model,
            provider=provider
        )
        
        if not md_success:
            logger.error("Markdown 轉換失敗")
            success = False
    
    # 處理 PPT
    if output_format in ["pptx", "both"]:
        output_ppt = os.path.join(
            os.path.dirname(slides_folder), 
            "slides.pptx"
        )
        
        ppt_success = process_images_to_ppt(
            image_dir=slides_folder,
            output_ppt=output_ppt,
            title="投影片簡報",
            use_llm=(api_key is not None),
            api_key=api_key,
            model=model,
            provider=provider
        )
        
        if not ppt_success:
            logger.error("PowerPoint 轉換失敗")
            success = False
    
    return success


def main():
    """主函數，處理命令列引數並執行檔案處理"""
    parser = argparse.ArgumentParser(
        description="MarkItDown PPT 轉換工具 - 將圖片轉換為 Markdown 並生成 PPT"
    )
    
    # 必要引數
    parser.add_argument(
        "image_dir",
        help="圖片資料夾路徑"
    )
    parser.add_argument(
        "output_ppt",
        help="輸出 PPT 檔案路徑"
    )
    
    # 可選引數
    parser.add_argument(
        "--markdown", "-m",
        help="輸出 Markdown 檔案路徑，如果未指定則在 PPT 檔案旁創建"
    )
    parser.add_argument(
        "--title", "-t",
        default="圖片簡報",
        help="簡報標題"
    )
    parser.add_argument(
        "--no-llm", "-n",
        action="store_true",
        help="不使用 LLM 處理圖片"
    )
    parser.add_argument(
        "--api-key", "-k",
        help="OpenAI API 金鑰，如果未指定則嘗試從環境變數獲取"
    )
    parser.add_argument(
        "--model", "-mo",
        default="gpt-4o-mini",
        help="使用的 OpenAI 模型，預設為 gpt-4o-mini"
    )
    parser.add_argument(
        "--template", "-tp",
        help="PPT 模板檔案路徑"
    )
    parser.add_argument(
        "--verbose", "-v",
        action="store_true",
        help="顯示詳細資訊"
    )
    
    args = parser.parse_args()
    
    # 設定日誌級別
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
    
    # 執行處理
    result = process_images_to_ppt(
        image_dir=args.image_dir,
        output_ppt=args.output_ppt,
        markdown_file=args.markdown,
        title=args.title,
        use_llm=not args.no_llm,
        api_key=args.api_key,
        model=args.model,
        template_file=args.template
    )
    
    # 顯示結果
    if result:
        logger.info("處理完成!")
        sys.exit(0)
    else:
        logger.error("處理失敗!")
        sys.exit(1)


if __name__ == "__main__":
    main() 