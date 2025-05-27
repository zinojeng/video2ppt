#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
瀏覽器捕獲工具：專門用於擷取瀏覽器中的投影片
"""

import os
import sys
import time
import argparse
import tkinter as tk
from tkinter import messagebox, Scale
import numpy as np
import cv2
import pyautogui
from datetime import datetime
from PIL import Image, ImageTk
from skimage.metrics import structural_similarity as ssim

class BrowserCapture:
    def __init__(self, root=None):
        # 創建主窗口
        if not root:
            self.root = tk.Tk()
            self.root.title("瀏覽器投影片捕獲工具")
            self.root.geometry("1000x700")
        else:
            self.root = root
        
        # 預設設置
        self.threshold = 0.95  # 相似度閾值
        self.interval = 0.5    # 檢測間隔（秒）
        self.output_folder = "slides"
        self.output_file = "slides.pptx"
        self.is_running = False
        self.is_paused = False
        self.last_frame = None
        self.slide_count = 0
        self.roi = None
        self.roi_selecting = False
        self.roi_start = None
        self.roi_end = None
        self.display_size = None
        
        # 確保輸出目錄存在
        if not os.path.exists(self.output_folder):
            os.makedirs(self.output_folder)
            
        # 設置界面
        self.setup_ui()
        
        # 顯示螢幕預覽
        self.show_screen_preview()
    
    def setup_ui(self):
        # 上方控制區域
        control_frame = tk.Frame(self.root)
        control_frame.pack(fill=tk.X, padx=10, pady=10)
        
        # 第一行 - 設置
        row1 = tk.Frame(control_frame)
        row1.pack(fill=tk.X, pady=5)
        
        tk.Label(row1, text="相似度閾值:").pack(side=tk.LEFT, padx=5)
        self.threshold_scale = Scale(row1, from_=0.5, to=1.0, resolution=0.01, 
                                    orient=tk.HORIZONTAL, length=200)
        self.threshold_scale.set(self.threshold)
        self.threshold_scale.pack(side=tk.LEFT, padx=5)
        
        tk.Label(row1, text="檢測間隔(秒):").pack(side=tk.LEFT, padx=5)
        self.interval_scale = Scale(row1, from_=0.1, to=5.0, resolution=0.1, 
                                   orient=tk.HORIZONTAL, length=200)
        self.interval_scale.set(self.interval)
        self.interval_scale.pack(side=tk.LEFT, padx=5)
        
        # 第二行 - 按鈕
        row2 = tk.Frame(control_frame)
        row2.pack(fill=tk.X, pady=5)
        
        self.start_btn = tk.Button(row2, text="開始捕獲", command=self.start_capture, 
                                  bg="#4CAF50", fg="white", width=10)
        self.start_btn.pack(side=tk.LEFT, padx=5)
        
        self.pause_btn = tk.Button(row2, text="暫停", command=self.pause_capture, 
                                  state=tk.DISABLED, width=10)
        self.pause_btn.pack(side=tk.LEFT, padx=5)
        
        self.stop_btn = tk.Button(row2, text="停止", command=self.stop_capture, 
                                 state=tk.DISABLED, width=10)
        self.stop_btn.pack(side=tk.LEFT, padx=5)
        
        self.reset_roi_btn = tk.Button(row2, text="重設區域", command=self.reset_roi, width=10)
        self.reset_roi_btn.pack(side=tk.LEFT, padx=5)
        
        # 右側狀態訊息
        self.status_var = tk.StringVar(value="就緒")
        status_label = tk.Label(row2, textvariable=self.status_var, 
                              bd=1, relief=tk.SUNKEN, anchor=tk.W)
        status_label.pack(side=tk.RIGHT, padx=5, fill=tk.X, expand=True)
        
        # 預覽區域
        self.canvas_frame = tk.Frame(self.root, bg="black")
        self.canvas_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        self.canvas = tk.Canvas(self.canvas_frame, bg="black")
        self.canvas.pack(fill=tk.BOTH, expand=True)
        
        # 綁定滑鼠事件
        self.canvas.bind("<ButtonPress-1>", self.on_mouse_down)
        self.canvas.bind("<B1-Motion>", self.on_mouse_move)
        self.canvas.bind("<ButtonRelease-1>", self.on_mouse_up)
        
        # 日誌區域
        log_frame = tk.Frame(self.root)
        log_frame.pack(fill=tk.X, padx=10, pady=5)
        
        tk.Label(log_frame, text="執行日誌:").pack(anchor=tk.W)
        self.log_text = tk.Text(log_frame, height=6, width=50)
        self.log_text.pack(fill=tk.X, expand=True)
        
        scrollbar = tk.Scrollbar(self.log_text)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.log_text.config(yscrollcommand=scrollbar.set)
        scrollbar.config(command=self.log_text.yview)
        
        self.log("程序已啟動")
        self.log("請框選要監控的投影片區域")
    
    def log(self, message):
        timestamp = datetime.now().strftime("%H:%M:%S")
        self.log_text.insert(tk.END, f"[{timestamp}] {message}\n")
        self.log_text.see(tk.END)
    
    def show_screen_preview(self):
        """顯示螢幕預覽"""
        try:
            # 獲取螢幕截圖
            screenshot = pyautogui.screenshot()
            
            # 轉換為 OpenCV 格式
            frame = np.array(screenshot)
            frame = cv2.cvtColor(frame, cv2.COLOR_RGB2BGR)
            
            # 顯示在畫布上
            self.display_frame(frame)
            
            self.log("已顯示螢幕預覽，請框選要監控的投影片區域")
            
            # 設置幀大小，用於ROI計算
            self.frame_size = (frame.shape[1], frame.shape[0])
        except Exception as e:
            self.log(f"無法獲取螢幕預覽: {str(e)}")
            messagebox.showerror("錯誤", f"無法獲取螢幕預覽: {str(e)}")
    
    def on_mouse_down(self, event):
        if not self.is_running and not self.is_paused:
            self.roi_selecting = True
            self.roi_start = (event.x, event.y)
            self.canvas.delete("roi")
    
    def on_mouse_move(self, event):
        if self.roi_selecting:
            self.roi_end = (event.x, event.y)
            self.canvas.delete("roi_temp")
            self.canvas.create_rectangle(
                self.roi_start[0], self.roi_start[1],
                self.roi_end[0], self.roi_end[1],
                outline="red", width=2, tags="roi_temp"
            )
    
    def on_mouse_up(self, event):
        if self.roi_selecting:
            self.roi_selecting = False
            self.roi_end = (event.x, event.y)
            
            # 計算 ROI
            x1, y1 = min(self.roi_start[0], self.roi_end[0]), min(self.roi_start[1], self.roi_end[1])
            x2, y2 = max(self.roi_start[0], self.roi_end[0]), max(self.roi_start[1], self.roi_end[1])
            
            # 保存實際像素範圍
            if hasattr(self, 'frame_size'):
                # 計算實際像素與顯示像素的比例
                ratio_x = self.frame_size[0] / self.display_size[0]
                ratio_y = self.frame_size[1] / self.display_size[1]
                
                # 計算實際像素位置
                self.roi = (
                    int(x1 * ratio_x), int(y1 * ratio_y),
                    int(x2 * ratio_x), int(y2 * ratio_y)
                )
            else:
                self.roi = (x1, y1, x2, y2)
            
            self.canvas.delete("roi_temp")
            self.canvas.create_rectangle(
                x1, y1, x2, y2,
                outline="green", width=2, tags="roi"
            )
            
            self.log(f"已選擇區域: ({x1},{y1}) - ({x2},{y2})")
    
    def reset_roi(self):
        self.roi = None
        self.canvas.delete("roi")
        self.log("已重設選擇區域")
        self.show_screen_preview()
    
    def display_frame(self, frame):
        """顯示幀到畫布上"""
        # 調整大小以適應畫布
        canvas_width = self.canvas.winfo_width()
        canvas_height = self.canvas.winfo_height()
        
        if canvas_width <= 1 or canvas_height <= 1:
            canvas_width = self.canvas_frame.winfo_width()
            canvas_height = self.canvas_frame.winfo_height()
        
        # 保持寬高比
        scale = min(canvas_width/frame.shape[1], canvas_height/frame.shape[0])
        new_width = int(frame.shape[1] * scale)
        new_height = int(frame.shape[0] * scale)
        
        self.display_size = (new_width, new_height)
        
        # 調整大小並轉換為 PIL 圖像
        frame_resized = cv2.resize(frame, (new_width, new_height))
        frame_rgb = cv2.cvtColor(frame_resized, cv2.COLOR_BGR2RGB)
        pil_img = Image.fromarray(frame_rgb)
        
        # 轉換為 Tkinter 可用的格式
        img_tk = ImageTk.PhotoImage(image=pil_img)
        
        # 保持引用以防垃圾回收
        self.current_image = img_tk
        
        # 更新畫布
        self.canvas.delete("frame")
        self.canvas.create_image(
            canvas_width//2, canvas_height//2,
            image=img_tk, anchor=tk.CENTER, tags="frame"
        )
        
        # 如果有 ROI，重新繪製
        if self.roi and not self.roi_selecting:
            x1, y1, x2, y2 = self.roi
            
            # 將實際像素位置轉換為顯示像素位置
            ratio_x = new_width / self.frame_size[0]
            ratio_y = new_height / self.frame_size[1]
            display_x1 = int(x1 * ratio_x)
            display_y1 = int(y1 * ratio_y)
            display_x2 = int(x2 * ratio_x)
            display_y2 = int(y2 * ratio_y)
            
            # 計算偏移量以居中顯示
            offset_x = (canvas_width - new_width) // 2
            offset_y = (canvas_height - new_height) // 2
            
            self.canvas.delete("roi")
            self.canvas.create_rectangle(
                display_x1 + offset_x, display_y1 + offset_y,
                display_x2 + offset_x, display_y2 + offset_y,
                outline="green", width=2, tags="roi"
            )
    
    def start_capture(self):
        """開始捕獲"""
        # 確保已選擇 ROI
        if not self.roi:
            messagebox.showerror("錯誤", "請先框選要監控的投影片區域")
            return
        
        # 更新設置
        self.threshold = float(self.threshold_scale.get())
        self.interval = float(self.interval_scale.get())
        
        # 更新界面狀態
        self.start_btn.config(state=tk.DISABLED)
        self.pause_btn.config(state=tk.NORMAL)
        self.stop_btn.config(state=tk.NORMAL)
        self.is_running = True
        self.is_paused = False
        
        self.log(f"開始捕獲")
        self.log(f"相似度閾值: {self.threshold}, 間隔: {self.interval}秒")
        
        # 啟動捕獲線程
        self.root.after(100, self.capture_loop)
    
    def capture_loop(self):
        """捕獲循環"""
        if not self.is_running or self.is_paused:
            return
        
        try:
            # 獲取截圖
            x1, y1, x2, y2 = self.roi
            width = x2 - x1
            height = y2 - y1
            
            screenshot = pyautogui.screenshot(region=(x1, y1, width, height))
            frame = np.array(screenshot)
            frame = cv2.cvtColor(frame, cv2.COLOR_RGB2BGR)
            
            # 顯示在畫布上
            self.display_current_roi(frame)
            
            # 檢查是否有變化
            if self.last_frame is not None:
                # 轉換為灰度圖像以進行比較
                gray1 = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
                gray2 = cv2.cvtColor(self.last_frame, cv2.COLOR_BGR2GRAY)
                
                # 計算相似度
                similarity = ssim(gray1, gray2)
                
                # 更新狀態
                self.status_var.set(f"相似度: {similarity:.4f}")
                
                # 如果幀間差異足夠大，保存截圖
                if similarity < self.threshold:
                    self.save_slide(frame)
                    self.last_frame = frame.copy()
            else:
                self.last_frame = frame.copy()
                self.save_slide(frame)
                
        except Exception as e:
            self.log(f"捕獲出錯: {str(e)}")
        
        # 安排下一次捕獲
        self.root.after(int(self.interval * 1000), self.capture_loop)
    
    def display_current_roi(self, frame):
        """顯示當前捕獲的 ROI 區域"""
        # 在小視窗中顯示 ROI
        small_canvas = tk.Canvas(self.root, width=320, height=240, bg="black")
        small_canvas.place(x=10, y=self.root.winfo_height() - 250)
        
        # 調整大小
        h, w = frame.shape[:2]
        scale = min(320/w, 240/h)
        new_width = int(w * scale)
        new_height = int(h * scale)
        
        frame_resized = cv2.resize(frame, (new_width, new_height))
        frame_rgb = cv2.cvtColor(frame_resized, cv2.COLOR_BGR2RGB)
        pil_img = Image.fromarray(frame_rgb)
        img_tk = ImageTk.PhotoImage(image=pil_img)
        
        # 保持引用
        self.roi_image = img_tk
        
        # 顯示圖像
        small_canvas.create_image(
            320//2, 240//2,
            image=img_tk, anchor=tk.CENTER
        )
    
    def save_slide(self, frame):
        """保存投影片"""
        slide_filename = os.path.join(self.output_folder, f"slide_{self.slide_count:03d}.png")
        cv2.imwrite(slide_filename, frame)
        
        self.slide_count += 1
        self.log(f"保存投影片 {self.slide_count}: {slide_filename}")
    
    def pause_capture(self):
        """暫停/繼續捕獲"""
        if self.is_running:
            if self.is_paused:
                self.is_paused = False
                self.pause_btn.config(text="暫停")
                self.log("繼續捕獲")
                self.capture_loop()
            else:
                self.is_paused = True
                self.pause_btn.config(text="繼續")
                self.log("暫停捕獲")
    
    def stop_capture(self):
        """停止捕獲"""
        self.is_running = False
        self.is_paused = False
        
        self.start_btn.config(state=tk.NORMAL)
        self.pause_btn.config(state=tk.DISABLED)
        self.stop_btn.config(state=tk.DISABLED)
        
        self.log("停止捕獲")
        
        if self.slide_count > 0:
            self.log(f"總共捕獲 {self.slide_count} 張投影片")
            
            # 詢問是否生成 PowerPoint
            if messagebox.askyesno("完成", 
                               f"已捕獲 {self.slide_count} 張投影片。\n是否立即生成 PowerPoint 文件?"):
                self.generate_ppt()
    
    def generate_ppt(self):
        """生成 PowerPoint 文件"""
        try:
            # 嘗試導入 python-pptx
            from pptx import Presentation
            from pptx.util import Inches
            
            # 創建 PPT
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]  # 空白佈局
            
            # 為每張圖片創建幻燈片
            for i in range(self.slide_count):
                slide_filename = os.path.join(self.output_folder, f"slide_{i:03d}.png")
                
                if os.path.exists(slide_filename):
                    slide = prs.slides.add_slide(slide_layout)
                    # 添加圖片並調整大小以適應幻燈片
                    slide.shapes.add_picture(slide_filename, Inches(0), Inches(0), 
                                           width=prs.slide_width, height=prs.slide_height)
            
            # 保存 PPT
            prs.save(self.output_file)
            self.log(f"已成功生成 PowerPoint 文件: {self.output_file}")
            messagebox.showinfo("成功", f"已成功生成 PowerPoint 文件: {self.output_file}")
            
        except ImportError:
            self.log("未找到 python-pptx 庫，請安裝：pip install python-pptx")
            messagebox.showerror("錯誤", "未找到 python-pptx 庫，請安裝：pip install python-pptx")
        except Exception as e:
            self.log(f"生成 PowerPoint 時出錯: {str(e)}")
            messagebox.showerror("錯誤", f"生成 PowerPoint 時出錯: {str(e)}")
    
    def run(self):
        """運行應用程序"""
        self.root.protocol("WM_DELETE_WINDOW", self.on_closing)
        self.root.mainloop()
    
    def on_closing(self):
        """窗口關閉處理"""
        if self.is_running:
            if messagebox.askyesno("確認", "捕獲仍在進行中。確定要退出嗎?"):
                self.stop_capture()
                self.root.destroy()
        else:
            self.root.destroy()

# 主函數
def main():
    parser = argparse.ArgumentParser(description="瀏覽器投影片捕獲工具")
    parser.add_argument("--threshold", type=float, default=0.95, help="相似度閾值 (0.5-1.0)")
    parser.add_argument("--interval", type=float, default=0.5, help="檢測間隔(秒)")
    parser.add_argument("--output", type=str, default="slides.pptx", help="輸出 PowerPoint 文件名")
    args = parser.parse_args()
    
    app = BrowserCapture()
    app.threshold = args.threshold
    app.threshold_scale.set(args.threshold)
    app.interval = args.interval
    app.interval_scale.set(args.interval)
    app.output_file = args.output
    
    app.run()

if __name__ == "__main__":
    main() 